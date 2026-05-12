"""Business Model + Industry slide redesign.

Single-slide pptx. Matches existing deck typography & palette.
Left third: compact triangulation visual. Right two-thirds: 4 styled
industry-signal cards. Bottom: navy convergence strip.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============ TOKENS (matched to existing deck) ============
NAVY        = RGBColor(0x0E, 0x1E, 0x3D)
NAVY_DEEP   = RGBColor(0x1F, 0x3A, 0x5F)
ORANGE      = RGBColor(0xE6, 0x6E, 0x2C)
ORANGE_HOT  = RGBColor(0xC8, 0x44, 0x1A)
ORANGE_NUM  = RGBColor(0xE1, 0x79, 0x3C)
ORANGE_BG   = RGBColor(0xFF, 0xF4, 0xEF)
ORANGE_BORD = RGBColor(0xE8, 0x90, 0x6A)
GREY_BG     = RGBColor(0xF3, 0xF4, 0xF6)
GREY_BORD   = RGBColor(0x6B, 0x72, 0x80)
GREY_BORD2  = RGBColor(0xE8, 0xE3, 0xDC)
TEXT_DARK   = RGBColor(0x11, 0x18, 0x27)
TEXT_GREY   = RGBColor(0x6B, 0x72, 0x80)
TEXT_GREY2  = RGBColor(0x70, 0x70, 0x70)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREEN       = RGBColor(0x2C, 0x7A, 0x4E)
BLUE_HW     = RGBColor(0xC8, 0x44, 0x1A)   # use orange-hot for hardware vertex (matches current)
BLUE_SUB    = RGBColor(0x1F, 0x3A, 0x5F)   # navy for subscription
GREEN_HEALTH= RGBColor(0x2C, 0x7A, 0x4E)   # green for health-data

FONT = "Calibri"
FONT_DISPLAY = "Georgia"

# ============ HELPERS ============
def add_textbox(slide, left, top, width, height, text, *,
                font=FONT, size=10, bold=False, italic=False, color=TEXT_DARK,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02);  tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = font; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_rich_textbox(slide, left, top, width, height, runs, *,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02);  tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    for txt, attrs in runs:
        run = p.add_run(); run.text = txt
        f = run.font
        f.name = attrs.get('font', FONT)
        f.size = Pt(attrs.get('size', 10))
        f.bold = attrs.get('bold', False)
        f.italic = attrs.get('italic', False)
        if 'color' in attrs:
            f.color.rgb = attrs['color']
    return tb

def add_rect(slide, left, top, width, height, *,
             fill=WHITE, line=GREY_BORD, line_width=0.5, no_line=False,
             shape=MSO_SHAPE.RECTANGLE):
    rect = slide.shapes.add_shape(shape,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    if fill is None:
        rect.fill.background()
    else:
        rect.fill.solid(); rect.fill.fore_color.rgb = fill
    if no_line:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line; rect.line.width = Pt(line_width)
    rect.shadow.inherit = False
    rect.text_frame.text = ""
    return rect

def add_line(slide, x1, y1, x2, y2, *, color=NAVY_DEEP, width=0.75):
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1),
                                     Inches(x2), Inches(y2))
    ln.line.color.rgb = color; ln.line.width = Pt(width)
    return ln

def add_circle(slide, cx, cy, radius, *, fill, line=None, line_width=1.0):
    """cx, cy are CENTER coords."""
    left = cx - radius; top = cy - radius
    d = radius * 2
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(left), Inches(top),
                                 Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(line_width)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh

# ============ BUILD ============
prs = Presentation()
prs.slide_width = Inches(20.0)
prs.slide_height = Inches(11.25)

BLANK = prs.slide_layouts[6]
slide = prs.slides.add_slide(BLANK)

# ---- TOP CHROME ----
add_textbox(slide, 0.50, 0.35, 19.00, 0.30,
            "WHOOP BUSINESS MODEL", size=11, bold=True, color=ORANGE)

# Title with orange emphasis
title_runs = [
    ("Three", {'font': FONT_DISPLAY, 'size': 34, 'bold': True, 'color': ORANGE_NUM}),
    (" companies in one  ·  the industry is converging on all three.",
     {'font': FONT_DISPLAY, 'size': 34, 'bold': True, 'color': NAVY}),
]
add_rich_textbox(slide, 0.50, 0.70, 19.00, 0.70, title_runs)

# Subtitle
add_textbox(slide, 0.50, 1.45, 19.00, 0.40,
            "WHOOP straddles consumer hardware, consumer subscription, and health-data platform — and four dated industry signals show the rest of the world moving toward that same intersection.",
            size=13, color=TEXT_GREY)

# Divider
add_line(slide, 0.50, 1.95, 19.50, 1.95, color=NAVY_DEEP, width=0.75)

# Footer
add_textbox(slide, 0.50, 10.78, 17.80, 0.22,
            "Sources: Grand View Research / Precedence Research (wearables TAM) · IDC (smart-ring share) · Dexcom 10-K (Stelo OTC Aug 2024) · WHOOP press (FDA ECG Apr 2025; CMS Innovation Center Apr 13, 2026) · FDA Makary Wellness guidance (Jan 2026) · iRhythm 10-K (CMS CPT codes 2023)",
            size=8.5, color=TEXT_GREY2)
add_textbox(slide, 18.50, 10.75, 1.20, 0.25,
            "03 / 17", size=9, color=TEXT_GREY, align=PP_ALIGN.RIGHT)

# ===========================================================================
# LEFT — TRIANGULATION VISUAL
# ===========================================================================
LEFT_X = 0.50
LEFT_W = 7.20

# Section label
add_textbox(slide, LEFT_X, 2.20, LEFT_W, 0.30,
            "§1  WHAT WHOOP IS  ·  NO PUBLIC COMP HOLDS IT TOGETHER",
            size=11, bold=True, color=ORANGE)

# Container box (subtle)
add_rect(slide, LEFT_X, 2.60, LEFT_W, 6.90, fill=WHITE, line=GREY_BORD2,
         line_width=0.75)

# Triangle geometry — equilateral-ish in the panel
# Panel center horizontally ~ 0.50 + 7.20/2 = 4.10
cx_center = LEFT_X + LEFT_W/2  # 4.10
top_y = 3.30   # top vertex y
bot_y = 7.40   # bottom vertices y

# Top vertex (Consumer Hardware)
top_cx = cx_center
top_cy = top_y
# Bottom-left vertex (Consumer Subscription)
bl_cx = LEFT_X + 1.10
bl_cy = bot_y
# Bottom-right vertex (Health Data Platform)
br_cx = LEFT_X + LEFT_W - 1.10
br_cy = bot_y

# Lines first (so they sit behind circles)
add_line(slide, top_cx, top_cy, bl_cx, bl_cy, color=NAVY, width=1.2)
add_line(slide, top_cx, top_cy, br_cx, br_cy, color=NAVY, width=1.2)
add_line(slide, bl_cx, bl_cy, br_cx, br_cy, color=NAVY, width=1.2)

# Vertex circles
add_circle(slide, top_cx, top_cy, 0.15, fill=BLUE_HW)
add_circle(slide, bl_cx,  bl_cy,  0.15, fill=BLUE_SUB)
add_circle(slide, br_cx,  br_cy,  0.15, fill=GREEN_HEALTH)

# Vertex labels — Consumer Hardware (above top vertex)
add_textbox(slide, top_cx - 1.50, top_cy - 0.65, 3.00, 0.30,
            "Consumer Hardware",
            size=13, bold=True, color=BLUE_HW, align=PP_ALIGN.CENTER)
add_textbox(slide, top_cx - 1.50, top_cy - 0.36, 3.00, 0.24,
            "GARMIN  ·  FITBIT  ·  GOPRO",
            size=9, bold=True, color=TEXT_GREY, align=PP_ALIGN.CENTER)

# Bottom-left — Consumer Subscription
add_textbox(slide, bl_cx - 1.30, bl_cy + 0.20, 2.60, 0.30,
            "Consumer Subscription",
            size=13, bold=True, color=BLUE_SUB, align=PP_ALIGN.CENTER)
add_textbox(slide, bl_cx - 1.30, bl_cy + 0.50, 2.60, 0.24,
            "PELOTON  ·  SONOS  ·  ROKU",
            size=9, bold=True, color=TEXT_GREY, align=PP_ALIGN.CENTER)

# Bottom-right — Health Data Platform
add_textbox(slide, br_cx - 1.30, br_cy + 0.20, 2.60, 0.30,
            "Health Data Platform",
            size=13, bold=True, color=GREEN_HEALTH, align=PP_ALIGN.CENTER)
add_textbox(slide, br_cx - 1.30, br_cy + 0.50, 2.60, 0.24,
            "DEXCOM  ·  MASIMO  ·  IRHYTHM",
            size=9, bold=True, color=TEXT_GREY, align=PP_ALIGN.CENTER)

# Center — WHOOP pill
pill_cx = cx_center
pill_cy = (top_cy + bot_y) / 2 + 0.10  # slight downward bias for visual weight
pill_w = 2.60; pill_h = 0.55
add_rect(slide, pill_cx - pill_w/2, pill_cy - pill_h/2, pill_w, pill_h,
         fill=WHITE, line=ORANGE_HOT, line_width=2.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, pill_cx - pill_w/2, pill_cy - 0.15, pill_w, 0.30,
            "WHOOP — AT THE INTERSECTION",
            size=10, bold=True, color=ORANGE_HOT, align=PP_ALIGN.CENTER)

# Caption below pill
add_textbox(slide, LEFT_X + 0.40, 8.50, LEFT_W - 0.80, 0.30,
            "Series G prices the intersection, not any single bucket.",
            size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_textbox(slide, LEFT_X + 0.40, 8.80, LEFT_W - 0.80, 0.60,
            "~19× recognized rev / ~9× bookings — above every public tier (highest single bucket: 6.7× hardware, 5.96× health-data). The premium IS the four-vector combination no listed company holds.",
            size=10, color=TEXT_GREY, align=PP_ALIGN.CENTER)

# ===========================================================================
# RIGHT — INDUSTRY SIGNALS (4 cards)
# ===========================================================================
RIGHT_X = 8.10
RIGHT_W = 11.40

# Section label
add_textbox(slide, RIGHT_X, 2.20, RIGHT_W, 0.30,
            "§2  WHERE THE INDUSTRY IS MOVING  ·  4 DATED SIGNALS",
            size=11, bold=True, color=ORANGE)

# 4 signal cards, evenly distributed
card_top = 2.60
card_h = 1.55
card_gap = 0.18
card_w = RIGHT_W

signals = [
    {
        "num": "01",
        "headline": "Wearable market doubling by 2033.",
        "stat": "$87B → $230-700B  ·  12-14% CAGR  ·  611M units shipped 2025",
        "source": "Grand View Research · Precedence Research · IDC",
        "whoop_read": "The rising tide. Even flat share gain compounds — WHOOP doesn't need to take share to clear its DCF.",
    },
    {
        "num": "02",
        "headline": "Medical-grade and consumer wearables are converging — from both directions.",
        "stat": "Dexcom Stelo OTC (Aug 2024)  ←→  WHOOP FDA-cleared ECG (Apr 2025) · BPI (Jan 2026)",
        "source": "Dexcom 10-K · WHOOP press",
        "whoop_read": "The dominant industry vector. WHOOP's healthcare pivot isn't an outlier bet — it's the direction of travel.",
    },
    {
        "num": "03",
        "headline": "Regulatory tailwind: FDA expanded wellness exemptions.",
        "stat": "Jan 2026 Makary guidance  ·  Optical BP measurement permitted in wellness category, no 510(k) required",
        "source": "FDA Commissioner Makary guidance · MedTech Dive",
        "whoop_read": "Lower compliance cost for sensor features. Validates WHOOP's BPI relaunch path; lifts attach potential on Healthspan.",
    },
    {
        "num": "04",
        "headline": "CMS reimbursement gates are opening to wearable platforms.",
        "stat": "iRhythm: regional MAC → national CPT codes  ·  WHOOP at CMS Innovation Center (Apr 13, 2026)",
        "source": "iRhythm 10-K (CPT 93243/93247, $213/$224)  ·  WHOOP / CMS press",
        "whoop_read": "Path from consumer-paid $30/mo to payer-paid $200+/episode is precedented. The Bull case lives here.",
    },
]

for i, sig in enumerate(signals):
    y = card_top + i * (card_h + card_gap)

    # Card body
    add_rect(slide, RIGHT_X, y, card_w, card_h,
             fill=WHITE, line=GREY_BORD2, line_width=0.75)

    # Left index bar (orange tab)
    add_rect(slide, RIGHT_X, y, 0.85, card_h,
             fill=ORANGE_BG, line=ORANGE_BG, no_line=True)
    add_textbox(slide, RIGHT_X, y + 0.30, 0.85, 0.90,
                sig["num"], font=FONT_DISPLAY, size=36, bold=True,
                color=ORANGE_HOT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Right content
    cx = RIGHT_X + 1.00
    cw = card_w - 1.10

    # Headline
    add_textbox(slide, cx, y + 0.08, cw, 0.32,
                sig["headline"], size=13, bold=True, color=NAVY)

    # Stat / evidence line
    add_textbox(slide, cx, y + 0.42, cw, 0.30,
                sig["stat"], size=10.5, color=TEXT_DARK)

    # Source (small italic)
    add_textbox(slide, cx, y + 0.72, cw, 0.22,
                sig["source"], size=8.5, italic=True, color=TEXT_GREY)

    # Divider hairline
    add_line(slide, cx, y + 0.97, cx + cw - 0.05, y + 0.97,
             color=GREY_BORD2, width=0.5)

    # WHOOP read (orange label + body)
    add_rich_textbox(slide, cx, y + 1.05, cw, 0.42, [
        ("WHOOP read  ·  ", {'size': 9, 'bold': True, 'color': ORANGE_HOT}),
        (sig["whoop_read"], {'size': 10, 'color': TEXT_DARK}),
    ])

# ===========================================================================
# BOTTOM CONVERGENCE STRIP
# ===========================================================================
strip_top = 9.65
strip_h = 0.95
add_rect(slide, 0.50, strip_top, 19.00, strip_h,
         fill=NAVY_DEEP, line=NAVY_DEEP, no_line=True)

# Orange leader
add_textbox(slide, 0.70, strip_top + 0.14, 18.60, 0.26,
            "THE CONVERGENCE",
            size=10, bold=True, color=ORANGE)

# Body line
convergence_runs = [
    ("Every industry signal moves toward ", {'size': 13, 'color': WHITE}),
    ("the same intersection WHOOP already occupies", {'size': 13, 'bold': True, 'color': ORANGE}),
    (".  Bull case 30% weight = P(all 4 vectors activate) — not a hopeful prior, an industry-aligned one.",
     {'size': 13, 'color': WHITE}),
]
add_rich_textbox(slide, 0.70, strip_top + 0.42, 18.60, 0.45, convergence_runs)

# ============ SAVE ============
OUT = "/home/user/whoop_val/Whoop Pitch - BM redesign.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
