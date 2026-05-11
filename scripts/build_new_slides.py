"""Build 4 new slides matching existing deck format.

Appends:
  - IPO Mechanics (replaces nothing; new slide)
  - Revenue / Members / ARPU drivers (replacement for slide 5)
  - Business Model (replacement for slide 3)
  - Verdict (replacement for slide 14)

Output: Whoop Pitch - test.NEW.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
from copy import deepcopy

# ============= DESIGN TOKENS (extracted from existing deck) =============
NAVY        = RGBColor(0x0E, 0x1E, 0x3D)   # primary headline navy
NAVY_DEEP   = RGBColor(0x1F, 0x3A, 0x5F)   # rectangles / connectors
NAVY_DARK   = RGBColor(0x1C, 0x2B, 0x3A)   # accents
ORANGE      = RGBColor(0xE6, 0x6E, 0x2C)   # §-section labels
ORANGE_HOT  = RGBColor(0xC8, 0x44, 0x1A)   # featured/highlight
ORANGE_NUM  = RGBColor(0xE1, 0x79, 0x3C)   # key numbers in headline
ORANGE_BG   = RGBColor(0xFF, 0xF4, 0xEF)   # callout bg
ORANGE_BG2  = RGBColor(0xFF, 0xF8, 0xF5)   # subtle bg
ORANGE_LITE = RGBColor(0xF5, 0xD5, 0xC8)   # progress fill
ORANGE_BORD = RGBColor(0xE8, 0x90, 0x6A)   # callout border
GREY_BG     = RGBColor(0xF3, 0xF4, 0xF6)   # rectangle bg
GREY_BORD   = RGBColor(0x6B, 0x72, 0x80)   # rectangle border
GREY_BORD2  = RGBColor(0xE8, 0xE3, 0xDC)   # subtle border
GREY_BORD3  = RGBColor(0xDE, 0xDA, 0xD4)   # gridline
TEXT_DARK   = RGBColor(0x11, 0x18, 0x27)   # primary body
TEXT_GREY   = RGBColor(0x6B, 0x72, 0x80)   # secondary body
TEXT_GREY2  = RGBColor(0x70, 0x70, 0x70)   # source footer
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREEN       = RGBColor(0x2C, 0x7A, 0x4E)   # bull
RED         = RGBColor(0xB0, 0x3A, 0x2E)   # bear

FONT = "Calibri"
FONT_DISPLAY = "Georgia"

# ============= LAYOUT CONSTANTS =============
SLIDE_W = 20.0   # inches
SLIDE_H = 11.25

TOP_LABEL_T = 0.35
TITLE_T = 0.70
SUB_T = 1.45
DIVIDER_T = 1.85
CONTENT_T = 2.10
FOOTER_T = 10.78
PAGENUM_T = 10.75

# ============= HELPERS =============
def add_textbox(slide, left, top, width, height, text, *,
                font=FONT, size=10, bold=False, italic=False, color=TEXT_DARK,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_rich_textbox(slide, left, top, width, height, runs, *,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs: list of (text, dict-of-font-attrs)"""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for i, (txt, attrs) in enumerate(runs):
        run = p.add_run()
        run.text = txt
        f = run.font
        f.name = attrs.get('font', FONT)
        f.size = Pt(attrs.get('size', 10))
        f.bold = attrs.get('bold', False)
        f.italic = attrs.get('italic', False)
        if 'color' in attrs:
            f.color.rgb = attrs['color']
    return tb

def add_rect(slide, left, top, width, height, *,
             fill=WHITE, line=GREY_BORD, line_width=0.5, no_line=False):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    if fill is None:
        rect.fill.background()
    else:
        rect.fill.solid()
        rect.fill.fore_color.rgb = fill
    if no_line:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(line_width)
    rect.shadow.inherit = False
    # nuke default text body
    rect.text_frame.text = ""
    return rect

def add_line(slide, x1, y1, x2, y2, *, color=NAVY_DEEP, width=0.75):
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1),
                                     Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    return ln

def slide_chrome(slide, top_label, title, title_orange_run=None, subtitle=None,
                 page_label="X / 17", sources=""):
    """Standard top + bottom chrome matching slide 7."""
    # Top label
    add_textbox(slide, 0.50, TOP_LABEL_T, 19.00, 0.30,
                top_label, size=11, bold=True, color=ORANGE)
    # Title (Georgia 30 bold navy, optional orange middle run)
    if title_orange_run is None:
        add_textbox(slide, 0.50, TITLE_T, 19.00, 0.65,
                    title, font=FONT_DISPLAY, size=30, bold=True, color=NAVY)
    else:
        # title = (left, orange, right) tuple
        runs = [
            (title_orange_run[0], {'font': FONT_DISPLAY, 'size': 30, 'bold': True, 'color': NAVY}),
            (title_orange_run[1], {'font': FONT_DISPLAY, 'size': 30, 'bold': True, 'color': ORANGE_NUM}),
            (title_orange_run[2], {'font': FONT_DISPLAY, 'size': 30, 'bold': True, 'color': NAVY}),
        ]
        add_rich_textbox(slide, 0.50, TITLE_T, 19.00, 0.65, runs)
    # Subtitle
    if subtitle:
        add_textbox(slide, 0.50, SUB_T, 19.00, 0.40,
                    subtitle, size=13, color=TEXT_GREY)
    # Divider
    add_line(slide, 0.50, DIVIDER_T, 19.50, DIVIDER_T,
             color=NAVY_DEEP, width=0.75)
    # Footer
    add_textbox(slide, 0.50, FOOTER_T, 17.80, 0.22,
                sources, size=8.5, color=TEXT_GREY2)
    add_textbox(slide, 18.50, PAGENUM_T, 1.20, 0.25,
                page_label, size=9, color=TEXT_GREY,
                align=PP_ALIGN.RIGHT)

def section_label(slide, left, top, width, text):
    add_textbox(slide, left, top, width, 0.35,
                text, size=11, bold=True, color=ORANGE)

def table_header(slide, left, top, cols, *, size=9.5):
    """cols = list of (x_offset, width, text)"""
    for x_off, w, label in cols:
        add_textbox(slide, left + x_off, top, w, 0.25,
                    label, size=size, bold=True, color=TEXT_GREY)

def table_row(slide, left, top, cells, *, size=10, height=0.25,
              bold_idx=None, grey_idx=None):
    """cells = list of (x_offset, width, text [, align])"""
    for i, cell in enumerate(cells):
        if len(cell) == 4:
            x_off, w, txt, align = cell
        else:
            x_off, w, txt = cell
            align = PP_ALIGN.LEFT
        color = TEXT_DARK
        bold = False
        if grey_idx and i in grey_idx:
            color = TEXT_GREY
        if bold_idx and i in bold_idx:
            bold = True
        add_textbox(slide, left + x_off, top, w, height,
                    txt, size=size, bold=bold, color=color, align=align)

# ============= LOAD DECK =============
SOURCE = "/home/user/whoop_val/Whoop Pitch - test.pptx"
OUT = "/home/user/whoop_val/Whoop Pitch - test.NEW.pptx"

prs = Presentation(SOURCE)
BLANK = prs.slide_layouts[6]   # Blank layout (slide 7 uses this)

# ===========================================================================
# SLIDE A — IPO MECHANICS (new slide, follows IPO column on Slide 4)
# ===========================================================================
def build_ipo_slide(prs):
    slide = prs.slides.add_slide(BLANK)
    slide_chrome(
        slide,
        top_label="IMPLIED IPO RANGE — REFERENCE-CLASS STEP-UP",
        title=None,
        title_orange_run=("Series G × ", "1.45×", " empirical step-up → $14.69B Base IPO mark."),
        subtitle="Six consumer / health platform IPOs replace the assumed 25% IRR. Median step-up 1.45×; median annualized IRR 24.7% — anchored, not assumed.",
        page_label="05 / 17",
        sources="Sources: Pitchbook, S&P Capital IQ, SEC S-1 filings · Peloton (Sep 2019), Fitbit (Jun 2015), Oscar Health (Mar 2021), GoPro (Jun 2014), Hims & Hers (Jan 2021), Allbirds (Nov 2021)"
    )

    # ============ §1 REFERENCE CLASS TABLE ============
    section_label(slide, 0.50, 2.10, 12.00, "§1  REFERENCE CLASS  ·  6 LAST-PRIVATE-ROUND → IPO TRANSITIONS")
    add_rect(slide, 0.50, 2.50, 12.00, 4.05, fill=WHITE, line=GREY_BORD)

    # Header row
    cols = [
        (0.20, 1.40, "Company"),
        (1.65, 1.30, "Last priv. round"),
        (2.95, 0.90, "Round date"),
        (3.85, 0.85, "Post-$ ($B)"),
        (4.75, 0.85, "IPO date"),
        (5.65, 0.85, "IPO val. ($B)"),
        (6.55, 0.60, "Mo."),
        (7.20, 0.80, "Step-up"),
        (8.05, 0.80, "Ann. IRR"),
        (8.90, 3.00, "Commentary"),
    ]
    table_header(slide, 0.50, 2.62, cols)
    add_line(slide, 0.60, 2.92, 12.40, 2.92, color=GREY_BORD2, width=0.5)

    rows = [
        ("Peloton",        "Series F", "Aug 2018", "4.15",  "Sep 2019", "8.10",  "13",  "1.95×", "85.4%",  "Connected-fitness sub — closest WHOOP analog"),
        ("Fitbit",         "Series F", "Aug 2013", "0.30",  "Jun 2015", "4.10",  "22",  "13.7×", "316%",   "OUTLIER — pre-platform private round; excluded from median"),
        ("Oscar Health",   "Series G", "Aug 2018", "3.20",  "Mar 2021", "7.92",  "31",  "2.48×", "42.0%",  "Digital health / insurtech — same Series G stage"),
        ("GoPro",          "Foxconn",  "Dec 2012", "2.25",  "Jun 2014", "2.95",  "18",  "1.31×", "19.8%",  "Consumer wearable hardware — closest B1 analog"),
        ("Hims & Hers",    "Series C", "Jan 2019", "1.10",  "Jan 2021", "1.60",  "24",  "1.45×", "20.6%",  "Consumer health subscription (SPAC vehicle)"),
        ("Allbirds",       "Series F", "Sep 2020", "1.70",  "Nov 2021", "2.20",  "14",  "1.29×", "24.7%",  "DTC consumer brand"),
    ]
    y = 3.00
    for i, r in enumerate(rows):
        is_outlier = (i == 1)
        for ci, (txt, col) in enumerate(zip(r, cols)):
            x_off, w, _ = col
            color = TEXT_DARK
            bold = False
            if is_outlier:
                color = TEXT_GREY
            if ci in (7, 8) and not is_outlier:
                bold = True
            if ci == 0:
                bold = True
            align = PP_ALIGN.LEFT if ci in (0, 1, 9) else PP_ALIGN.RIGHT
            if ci == 2 or ci == 4:
                align = PP_ALIGN.LEFT
            add_textbox(slide, 0.50 + x_off, y, w, 0.27,
                        txt, size=9.5, bold=bold, color=color, align=align)
        y += 0.31

    # Footer notes inside §1 box
    add_textbox(slide, 0.70, 5.10, 11.50, 0.50,
                "Note: SPACs (Hims) and traditional IPOs treated equivalently — both represent first liquid public price. Median time-to-IPO ex-Fitbit ≈ 18 months.",
                size=8.5, color=TEXT_GREY, font=FONT)

    # Distribution stats sub-section inside same area
    add_line(slide, 0.70, 5.55, 12.30, 5.55, color=GREY_BORD2, width=0.5)
    add_textbox(slide, 0.70, 5.65, 4.00, 0.30,
                "Distribution (ex-Fitbit, n = 5)", size=10, bold=True, color=TEXT_DARK)
    stats = [
        ("25th percentile", "1.31×",  "BEAR"),
        ("Median",          "1.45×",  "BASE"),
        ("75th percentile", "1.95×",  "BULL"),
        ("Mean",            "1.69×",  "—"),
    ]
    sx = 0.70
    sy = 5.95
    for i, (label, val, tag) in enumerate(stats):
        add_textbox(slide, sx + i*2.85, sy,      1.70, 0.25, label,
                    size=9, color=TEXT_GREY)
        add_textbox(slide, sx + i*2.85, sy+0.22, 1.70, 0.30, val,
                    size=15, bold=True, color=TEXT_DARK)
        if tag != "—":
            tag_color = RED if tag == "BEAR" else (NAVY if tag == "BASE" else GREEN)
            add_textbox(slide, sx + i*2.85 + 0.95, sy+0.30, 0.85, 0.22, tag,
                        size=8.5, bold=True, color=tag_color)

    # ============ §2 APPLY STEP-UPS TO SERIES G ============
    section_label(slide, 12.80, 2.10, 6.70, "§2  APPLY TO SERIES G — IMPLIED IPO ENTERPRISE VALUE")
    add_rect(slide, 12.80, 2.50, 6.70, 2.65, fill=WHITE, line=GREY_BORD)

    # Header
    add_textbox(slide, 13.00, 2.62, 1.80, 0.25, "Scenario",
                size=9.5, bold=True, color=TEXT_GREY)
    add_textbox(slide, 14.80, 2.62, 0.90, 0.25, "Step-up",
                size=9.5, bold=True, color=TEXT_GREY, align=PP_ALIGN.RIGHT)
    add_textbox(slide, 15.80, 2.62, 1.20, 0.25, "Anchor",
                size=9.5, bold=True, color=TEXT_GREY, align=PP_ALIGN.RIGHT)
    add_textbox(slide, 17.20, 2.62, 2.10, 0.25, "Implied IPO EV ($B)",
                size=9.5, bold=True, color=TEXT_GREY, align=PP_ALIGN.RIGHT)
    add_line(slide, 12.90, 2.92, 19.40, 2.92, color=GREY_BORD2, width=0.5)

    rows2 = [
        ("BEAR",  "1.31×", "$10.10B", "$13.24B", RED),
        ("BASE",  "1.45×", "$10.10B", "$14.69B", NAVY),
        ("BULL",  "1.95×", "$10.10B", "$19.71B", GREEN),
    ]
    y = 3.05
    for tag, step, anchor, ev, color in rows2:
        add_textbox(slide, 13.00, y, 1.80, 0.32, tag,
                    size=11, bold=True, color=color)
        add_textbox(slide, 14.80, y, 0.90, 0.32, step,
                    size=11, color=TEXT_DARK, align=PP_ALIGN.RIGHT)
        add_textbox(slide, 15.80, y, 1.20, 0.32, anchor,
                    size=11, color=TEXT_GREY, align=PP_ALIGN.RIGHT)
        add_textbox(slide, 17.20, y, 2.10, 0.32, ev,
                    size=14, bold=True, color=color, align=PP_ALIGN.RIGHT)
        y += 0.46

    # Implied IRR row
    add_line(slide, 12.90, y+0.05, 19.40, y+0.05, color=GREY_BORD2, width=0.5)
    add_textbox(slide, 13.00, y+0.15, 4.00, 0.30,
                "Implied IRR (Series G entry → IPO, 18mo)",
                size=10, color=TEXT_GREY)
    add_textbox(slide, 17.20, y+0.15, 2.10, 0.30,
                "19.8% / 24.7% / 80.5%",
                size=10, bold=True, color=TEXT_DARK, align=PP_ALIGN.RIGHT)

    # ============ §3 CROSS-CHECK — IRR BACK-SOLVE ============
    section_label(slide, 12.80, 5.40, 6.70, "§3  CROSS-CHECK  ·  IS THE STEP-UP CONSISTENT WITH SERIES G IRR HURDLE?")
    add_rect(slide, 12.80, 5.80, 6.70, 1.40, fill=ORANGE_BG, line=ORANGE_BORD, line_width=0.75)

    add_textbox(slide, 13.00, 5.95, 6.40, 0.35,
                "Series G investors target ≥ 25% IRR over 12-18 mo to IPO.",
                size=10, color=TEXT_DARK)
    add_textbox(slide, 13.00, 6.25, 6.40, 0.50,
                "Empirical median (24.7% ann. IRR ex-Fitbit) lands inside that hurdle. The 1.45× step-up therefore is not an analyst dial — it is what comparable rounds have actually delivered. The IPO method validates Series G entry, not the other way around.",
                size=9.5, color=TEXT_DARK)
    add_textbox(slide, 13.00, 6.85, 6.40, 0.25,
                "Reference-class anchor replaces the prior 25% IRR back-solve at Assumptions C93:E93.",
                size=8.5, color=TEXT_GREY, font=FONT)

    # ============ §4 CAVEAT — INDEPENDENCE OF METHOD ============
    add_rect(slide, 0.50, 6.75, 12.00, 1.70, fill=NAVY_DEEP, line=NAVY_DEEP, no_line=True)
    add_textbox(slide, 0.70, 6.90, 11.60, 0.30,
                "WHY 15% WEIGHT — AND NOT MORE", size=11, bold=True, color=ORANGE)
    add_textbox(slide, 0.70, 7.20, 11.60, 1.20,
                "The Implied IPO method is forward-looking but not fully independent of Series G — it scales the $10.10B anchor by an empirically observed step-up multiplier. "
                "Two of the three reference-class anchors (1.31× / 1.95×) sit at percentile endpoints; the median (1.45×) is the live base case. "
                "Information overlap with Series G method means we cap the weight at 15% (vs. 40% Intrinsic, 35% Comps) — high enough to inject a forward anchor, low enough that the conclusion doesn't rest on it.",
                size=11, color=WHITE)

    # ============ §5 BOTTOM STRIP — FOOTBALL FIELD CONTRIBUTION ============
    add_rect(slide, 12.80, 7.30, 6.70, 1.15, fill=GREY_BG, line=GREY_BORD2)
    add_textbox(slide, 13.00, 7.42, 6.40, 0.30,
                "CONTRIBUTION TO WEIGHTED FOOTBALL FIELD",
                size=9.5, bold=True, color=TEXT_GREY)
    add_textbox(slide, 13.00, 7.72, 6.40, 0.30,
                "0.15  ×  $14.69B  =  $2.20B",
                size=15, bold=True, color=NAVY)
    add_textbox(slide, 13.00, 8.02, 6.40, 0.40,
                "Weighted avg without IPO method: $7.14B → +$2.20B → $9.09B blended.",
                size=9.5, color=TEXT_GREY)

    # ============ §6 BOTTOM BAR — KEY READS ============
    add_rect(slide, 0.50, 8.65, 19.00, 1.95, fill=WHITE, line=GREY_BORD2)
    add_textbox(slide, 0.70, 8.78, 18.60, 0.30,
                "WHAT THE EMPIRICAL STEP-UP TELLS US",
                size=11, bold=True, color=ORANGE)
    add_line(slide, 0.70, 9.12, 19.30, 9.12, color=GREY_BORD2, width=0.5)

    reads = [
        ("Median step-up exceeds 1.31× Bear floor.",
         "Even the bottom-quartile reference comp delivered +31% to its private mark by IPO. The Bear IPO case still clears Series G by +31%."),
        ("Time-to-IPO median ≈ 18 months.",
         "Aligns with Ahmed's '2027 IPO' guidance from Series G (Mar 2026). Suggests 12-18 mo discount window for any forward DCF anchor."),
        ("Connected-fitness analog (Peloton) delivered 1.95× — the 75th percentile.",
         "WHOOP's closest business-model analog at private→IPO transition justifies a Bull case meaningfully above median."),
        ("GoPro (1.31×) is the Bear floor.",
         "Wearable hardware without subscription compounding; if WHOOP fails the subscription thesis it converges to GoPro economics — and even then the IPO clears Series G."),
    ]
    cx = 0.70
    cy = 9.25
    cw = 4.55
    for i, (head, body) in enumerate(reads):
        add_textbox(slide, cx + i*cw, cy, cw - 0.10, 0.40,
                    head, size=10, bold=True, color=NAVY)
        add_textbox(slide, cx + i*cw, cy+0.42, cw - 0.10, 0.95,
                    body, size=9, color=TEXT_GREY)

    return slide


# ===========================================================================
# SLIDE B — REVENUE / MEMBERS / ARPU DRIVER SLIDE  (replaces Slide 5)
# ===========================================================================
def build_drivers_slide(prs):
    slide = prs.slides.add_slide(BLANK)
    slide_chrome(
        slide,
        top_label="DCF INPUTS — TWO-DRIVER FOUNDATION",
        title=None,
        title_orange_run=("Revenue = ", "Members × ARPU", ".  No add-on revenue lines."),
        subtitle="Members carries the analytical weight (4-lens triangulation → 10.9M Base 2033). ARPU is disciplined to inflation in Bear/Base; only Bull lifts — and only via 5 mechanism-anchored components.",
        page_label="06 / 17",
        sources="Sources: WHOOP Series G press · Sacra · Peloton 10-K (FY2025) · Grand View Research · Towards Healthcare · Damodaran · Internal cohort & ARPU build"
    )

    # ============ §1 IDENTITY STRIP ============
    add_rect(slide, 0.50, 2.10, 19.00, 1.55, fill=NAVY_DEEP, line=NAVY_DEEP, no_line=True)

    add_textbox(slide, 0.70, 2.22, 18.60, 0.28,
                "§1  THE IDENTITY  ·  REVENUE IS DERIVED, NOT FORECAST",
                size=11, bold=True, color=ORANGE)

    # Big formula
    formula_runs = [
        ("Revenue ($M)", {'font': FONT_DISPLAY, 'size': 20, 'bold': True, 'color': WHITE}),
        ("   =   ", {'font': FONT_DISPLAY, 'size': 20, 'color': ORANGE}),
        ("Avg Members (M)", {'font': FONT_DISPLAY, 'size': 20, 'bold': True, 'color': WHITE}),
        ("   ×   ", {'font': FONT_DISPLAY, 'size': 20, 'color': ORANGE}),
        ("Blended ARPU ($/yr)", {'font': FONT_DISPLAY, 'size': 20, 'bold': True, 'color': WHITE}),
    ]
    add_rich_textbox(slide, 0.70, 2.55, 18.60, 0.55, formula_runs)

    # Year endpoints
    year_rows = [
        ("2025 actual:",  "$521",   "1.9",  "$274",  TEXT_GREY),
        ("2033 Base:",    "$3,686", "10.6", "$347",  WHITE),
    ]
    yy = 3.15
    for label, rev, m, arpu, color in year_rows:
        add_textbox(slide, 0.70,  yy, 1.80, 0.26, label, size=11, color=TEXT_GREY)
        add_textbox(slide, 4.20,  yy, 2.40, 0.26, rev,   size=12, bold=True, color=color, align=PP_ALIGN.RIGHT)
        add_textbox(slide, 8.20,  yy, 2.40, 0.26, m,     size=12, bold=True, color=color, align=PP_ALIGN.RIGHT)
        add_textbox(slide, 12.40, yy, 3.00, 0.26, arpu,  size=12, bold=True, color=color, align=PP_ALIGN.RIGHT)
        yy += 0.24

    # Critical footnote
    add_textbox(slide, 0.70, 3.40, 18.60, 0.24,
                "Critical modeling rule: ARPU applies to AVERAGE members, NOT year-end. End-of-period × ARPU overstates revenue ~13% (the bookings-vs-revenue gap, R/B ratio 1.32× in 2025).",
                size=9, italic=True, color=ORANGE_LITE) if False else None
    add_rich_textbox(slide, 0.70, 3.40, 18.60, 0.24, [
        ("Critical rule:  ", {'size': 9, 'bold': True, 'color': ORANGE_LITE}),
        ("ARPU applies to AVERAGE members during the period, NOT year-end. End × ARPU overstates revenue ~13% — this is the bookings-vs-recognized-revenue gap (R/B ratio 1.32× in 2025, converging to 1.03× by 2033).",
         {'size': 9, 'italic': True, 'color': WHITE}),
    ])

    # ============ §2 MEMBERS PANEL (LEFT) ============
    section_label(slide, 0.50, 3.85, 9.40, "§2  MEMBERS  ·  4-LENS TRIANGULATION  →  10.9M BASE 2033")
    add_rect(slide, 0.50, 4.25, 9.40, 6.30, fill=WHITE, line=GREY_BORD)

    # 2A. Trajectory table
    add_textbox(slide, 0.70, 4.38, 9.00, 0.25,
                "Ending members by scenario (M)", size=9.5, bold=True, color=TEXT_GREY)
    add_line(slide, 0.70, 4.65, 9.70, 4.65, color=GREY_BORD2, width=0.5)

    years = ["2025", "2026", "2027", "2028", "2029", "2030", "2031", "2032", "2033"]
    bear  = ["2.5", "3.3",  "4.0",  "4.6",  "5.1",  "5.4",  "5.7",  "5.9",  "6.0"]
    base  = ["2.5", "3.7",  "5.1",  "6.4",  "7.6",  "8.6",  "9.6",  "10.3", "10.9"]
    bull  = ["2.5", "4.0",  "6.0",  "8.4",  "10.8", "12.7", "14.3", "15.6", "16.6"]

    # Year header
    add_textbox(slide, 0.70, 4.70, 1.10, 0.22, "", size=9, color=TEXT_GREY)
    for i, yr in enumerate(years):
        add_textbox(slide, 1.80 + i*0.88, 4.70, 0.86, 0.22, yr,
                    size=9, bold=True, color=TEXT_GREY, align=PP_ALIGN.RIGHT)

    for ri, (tag, data, color) in enumerate([("Bear", bear, RED), ("Base", base, NAVY), ("Bull", bull, GREEN)]):
        y = 4.95 + ri*0.28
        add_textbox(slide, 0.70, y, 1.10, 0.24, tag,
                    size=10, bold=True, color=color)
        for i, v in enumerate(data):
            is_endpoint = (i == len(data)-1)
            add_textbox(slide, 1.80 + i*0.88, y, 0.86, 0.24, v,
                        size=10, bold=is_endpoint, color=color if is_endpoint else TEXT_DARK,
                        align=PP_ALIGN.RIGHT)

    # 2B. Four-lens triangulation
    add_line(slide, 0.70, 5.85, 9.70, 5.85, color=GREY_BORD2, width=0.5)
    add_textbox(slide, 0.70, 5.95, 9.00, 0.25,
                "Four-lens triangulation — Base 2033 endpoint", size=9.5, bold=True, color=TEXT_GREY)

    lens = [
        ("Lens A · TAM × Share",       "25M TAM, 4% growth, 25% terminal share",                    "8.6M"),
        ("Lens B · Peloton stage-shift","WHOOP @ 2.5M ≈ Peloton @ IPO+2; 8yr forward adj for free HW","11.4M"),
        ("Lens C · Multi-comp median",  "Apple Watch shape · Oura penetration · Garmin trajectory",  "10.5M"),
        ("Lens D · Capital-led",        "$3.5B deployable S&M ÷ $245 all-in CAC − churn",            "8.4M"),
    ]
    ly = 6.25
    for label, justif, val in lens:
        add_textbox(slide, 0.70, ly, 3.20, 0.24, label, size=9, bold=True, color=TEXT_DARK)
        add_textbox(slide, 3.95, ly, 4.85, 0.24, justif, size=8.5, color=TEXT_GREY)
        add_textbox(slide, 8.85, ly, 0.85, 0.24, val, size=10, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
        ly += 0.23

    add_line(slide, 0.70, 7.20, 9.70, 7.20, color=GREY_BORD2, width=0.5)
    add_textbox(slide, 0.70, 7.28, 6.00, 0.28,
                "Median (4 lenses): 10.5M  →  Calibrated Base: 10.9M", size=11, bold=True, color=TEXT_DARK)
    add_textbox(slide, 0.70, 7.55, 8.50, 0.22,
                "Calibrated slightly above median to weight Lens C (strongest data quality).",
                size=8.5, color=TEXT_GREY)

    # 2C. Member sensitivity headline
    add_rect(slide, 0.70, 7.95, 9.00, 1.10, fill=ORANGE_BG, line=ORANGE_BORD, line_width=0.75)
    add_textbox(slide, 0.85, 8.05, 8.70, 0.25,
                "WHY MEMBERS DOMINATES THE TORNADO",
                size=10, bold=True, color=ORANGE_HOT)
    add_textbox(slide, 0.85, 8.30, 8.70, 0.65,
                "Members swing DCF Equity $7.5B → $15.8B Bear→Bull — ~3× the next-largest driver (terminal multiple). "
                "All analytical effort allocated proportionally: members get 4 independent lenses; ARPU stays disciplined.",
                size=10, color=TEXT_DARK)

    # 2D. Cohort churn note
    add_textbox(slide, 0.70, 9.15, 9.00, 0.22,
                "Cohort engine: Y1 22% / Y2 15% / Y3+ 12% churn → blended ~15% steady-state, ~17% during hypergrowth (new cohorts dominate).",
                size=8.5, color=TEXT_GREY, italic=True)
    add_textbox(slide, 0.70, 9.38, 9.00, 0.22,
                "Diagnostic guardrails (Members tab §9): LTV/CAC 4.3-5.4× ✓ · TAM penetration 41% terminal ✓ · S&M % rev within phase bands.",
                size=8.5, color=TEXT_GREY, italic=True)
    add_textbox(slide, 0.70, 9.60, 9.00, 0.22,
                "Bookings ≠ Revenue: $1.1B Series G headline = bookings run-rate (End-Members × ARPU); recognized rev (Avg × ARPU) = $521M.",
                size=8.5, color=TEXT_GREY, italic=True)

    # ============ §3 ARPU PANEL (RIGHT) ============
    section_label(slide, 10.10, 3.85, 9.40, "§3  ARPU  ·  ASYMMETRIC SCENARIO DISCIPLINE")
    add_rect(slide, 10.10, 4.25, 9.40, 6.30, fill=WHITE, line=GREY_BORD)

    # 3A. Scenario growth table
    add_textbox(slide, 10.30, 4.38, 9.00, 0.25,
                "ARPU growth by scenario (constant rate, compound)", size=9.5, bold=True, color=TEXT_GREY)
    add_line(slide, 10.30, 4.65, 19.30, 4.65, color=GREY_BORD2, width=0.5)

    arpu_hdr = [(0.00, 1.40, "Scenario"),
                (1.40, 1.40, "Growth"),
                (2.80, 1.60, "2025 ARPU"),
                (4.40, 1.60, "2033 ARPU"),
                (6.00, 1.40, "$/mo '33"),
                (7.40, 1.60, "Mechanism")]
    for x_off, w, lbl in arpu_hdr:
        align = PP_ALIGN.LEFT if x_off in (0.00, 7.40) else PP_ALIGN.RIGHT
        add_textbox(slide, 10.30 + x_off, 4.70, w, 0.22, lbl,
                    size=9, bold=True, color=TEXT_GREY, align=align)

    arpu_rows = [
        ("Bear", "2.5%", "$274", "$334", "$27.82", "Inflation − 50bps; consumer-only", RED),
        ("Base", "3.0%", "$274", "$347", "$28.92", "CPI baseline; no mix shift",        NAVY),
        ("Bull", "6.0%", "$274", "$437", "$36.39", "5-component decomp (below)",        GREEN),
    ]
    y = 4.95
    for tag, g, a25, a33, mo, mech, color in arpu_rows:
        add_textbox(slide, 10.30,        y, 1.40, 0.24, tag,    size=10, bold=True, color=color)
        add_textbox(slide, 10.30 + 1.40, y, 1.40, 0.24, g,      size=10, color=TEXT_DARK,  align=PP_ALIGN.RIGHT)
        add_textbox(slide, 10.30 + 2.80, y, 1.60, 0.24, a25,    size=10, color=TEXT_GREY,  align=PP_ALIGN.RIGHT)
        add_textbox(slide, 10.30 + 4.40, y, 1.60, 0.24, a33,    size=10, bold=True, color=color, align=PP_ALIGN.RIGHT)
        add_textbox(slide, 10.30 + 6.00, y, 1.40, 0.24, mo,     size=10, color=TEXT_DARK,  align=PP_ALIGN.RIGHT)
        add_textbox(slide, 10.30 + 7.40, y, 1.85, 0.24, mech,   size=9, color=TEXT_GREY)
        y += 0.28

    # 3B. Bull decomposition stacked
    add_line(slide, 10.30, 5.85, 19.30, 5.85, color=GREY_BORD2, width=0.5)
    add_textbox(slide, 10.30, 5.95, 9.00, 0.25,
                "Bull 6.0% = sum of 5 mechanism-anchored components (each tied to a commercial event)",
                size=9.5, bold=True, color=TEXT_GREY)

    bull = [
        ("Inflation baseline",                       "3.00%", "CPI+ baseline — consistent with Base"),
        ("Tier-mix shift to Peak / Life",            "1.00%", "Healthcare-positioning brand effect; Spotify Family analog"),
        ("Labs / BP / Healthspan attach",            "1.00%", "Premium feature monetization; BP launched Jan 2026"),
        ("B2B Unite share growth (20% → 30%)",       "0.30%", "Per-seat pricing 4-7× consumer"),
        ("Payer-paid revenue uplift",                "0.70%", "Direct reimbursement; loosest component (0.5-1.0% range)"),
    ]
    by = 6.25
    for comp, val, just in bull:
        add_textbox(slide, 10.30, by, 3.80, 0.24, comp, size=9, color=TEXT_DARK)
        add_textbox(slide, 14.20, by, 0.90, 0.24, val,  size=10, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
        add_textbox(slide, 15.20, by, 4.10, 0.24, just, size=8.5, color=TEXT_GREY)
        by += 0.22

    # Total
    add_line(slide, 14.20, 7.35, 15.10, 7.35, color=NAVY_DEEP, width=1)
    add_textbox(slide, 10.30, 7.40, 3.80, 0.26, "Bull Total",
                size=10, bold=True, color=TEXT_DARK)
    add_textbox(slide, 14.20, 7.40, 0.90, 0.26, "6.00%",
                size=11, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)

    # 3C. Reality check vs WHOOP tiers
    add_rect(slide, 10.30, 7.85, 9.00, 1.20, fill=ORANGE_BG, line=ORANGE_BORD, line_width=0.75)
    add_textbox(slide, 10.45, 7.95, 8.70, 0.25,
                "REALITY CHECK  ·  BULL $36.39/MO LANDS BETWEEN PEAK ($33) AND LIFE ($49)",
                size=10, bold=True, color=ORANGE_HOT)
    tier_ref = [
        ("WHOOP One",  "$30/mo",  "Entry tier · current floor",         TEXT_GREY),
        ("WHOOP Peak", "$33/mo",  "Bull endpoint sits 10% above this",  TEXT_DARK),
        ("WHOOP Life", "$49/mo",  "Already exceeds Bull endpoint",      TEXT_GREY),
    ]
    ty = 8.25
    for label, val, note, color in tier_ref:
        add_textbox(slide, 10.45, ty, 1.80, 0.24, label, size=9, bold=True, color=color)
        add_textbox(slide, 12.30, ty, 1.00, 0.24, val,   size=10, bold=True, color=color)
        add_textbox(slide, 13.35, ty, 5.80, 0.24, note,  size=9, color=TEXT_GREY)
        ty += 0.22

    # 3D. Why ARPU doesn't dominate the tornado
    add_textbox(slide, 10.30, 9.15, 9.00, 0.22,
                "Why ARPU growth is rank-5 in the tornado: range $334-$437 = $103 spread × 10.9M = $1.1B revenue swing. Members range = 6.0M-16.6M = $3.7B swing.",
                size=8.5, color=TEXT_GREY, italic=True)
    add_textbox(slide, 10.30, 9.38, 9.00, 0.22,
                "Discipline argument: Base case stays at inflation only. Bull lifts only because each mechanism corresponds to an observable commercial event (healthcare reimbursement clearance).",
                size=8.5, color=TEXT_GREY, italic=True)
    add_textbox(slide, 10.30, 9.60, 9.00, 0.22,
                "ARPU growth in Bear/Base is asymptotically capped at inflation. Bull = consequence of scenario, not independent dial. Single most-defended assumption in Q&A.",
                size=8.5, color=TEXT_GREY, italic=True)

    return slide


# ===========================================================================
# SLIDE C — BUSINESS MODEL (replaces Slide 3)
# ===========================================================================
def build_bm_slide(prs):
    slide = prs.slides.add_slide(BLANK)
    slide_chrome(
        slide,
        top_label="WHOOP BUSINESS MODEL",
        title=None,
        title_orange_run=("Four commercial vectors. ", "No public comp", " holds them together — that gap is the thesis."),
        subtitle="Today's monetization (left) plus four active expansion vectors (right) with dated evidence. The +$1.5B premium intrinsic captures over Series G IS the gap public-comp methodology cannot price.",
        page_label="03 / 17",
        sources="Sources: WHOOP Series G press · Sacra · Pitchbook · S&P Capital IQ · NYT DealBook · BMCT · BusinessWire · WHOOP press releases (CMS Innovation Center Apr 13, 2026 · Hiring surge · Advanced Labs)"
    )

    # ============ §1 LEFT — TODAY'S MONETIZATION ============
    section_label(slide, 0.50, 2.10, 9.40, "§1  WHAT THEY MONETIZE TODAY  ·  ~$521M RECOGNIZED FY25 / $1.1B BOOKINGS RUN-RATE")
    add_rect(slide, 0.50, 2.50, 9.40, 4.40, fill=WHITE, line=GREY_BORD)

    # Header
    hdr = [(0.00, 1.95, "Vector"),
           (1.95, 1.70, "What it is"),
           (3.65, 1.30, "2025 mix"),
           (4.95, 1.85, "Unit econ"),
           (6.80, 2.40, "Public analog")]
    for x_off, w, lbl in hdr:
        add_textbox(slide, 0.70 + x_off, 2.62, w, 0.25, lbl,
                    size=9.5, bold=True, color=TEXT_GREY)
    add_line(slide, 0.70, 2.92, 9.70, 2.92, color=GREY_BORD2, width=0.5)

    rows = [
        ("Consumer subscription", "$30/mo Pro tier (One/Peak/Life $199/$239/$359)",
         "~80% recurring", "LTV/CAC 4.3-5.4× · 78% sub GM",  "Spotify · Roku"),
        ("B2B Unite",              "Employer / team contracts (sports, military, healthcare orgs)",
         "~20% recurring",  "~$260 per-seat ARPU · 4-7× consumer", "No clean comp"),
        ("Health data / clinical", "FDA-cleared ECG · BPI · Healthspan · Advanced Labs",
         "Attach to sub",   "Premium feature monetization",  "Dexcom · iRhythm"),
        ("Hardware",               "Bundled-free strap (CAC line, not revenue)",
         "$0 revenue",      "BOM ~$50-80 · ~19-mo payback",   "Garmin (ref only)"),
    ]
    y = 3.00
    for vec, what, mix, econ, analog in rows:
        add_textbox(slide, 0.70,        y, 1.85, 0.65, vec,    size=10, bold=True,  color=NAVY)
        add_textbox(slide, 0.70 + 1.95, y, 1.65, 0.65, what,   size=9,  color=TEXT_DARK)
        add_textbox(slide, 0.70 + 3.65, y, 1.25, 0.65, mix,    size=9,  color=TEXT_GREY)
        add_textbox(slide, 0.70 + 4.95, y, 1.80, 0.65, econ,   size=9,  color=TEXT_DARK)
        add_textbox(slide, 0.70 + 6.80, y, 2.40, 0.65, analog, size=9,  color=TEXT_GREY)
        y += 0.95

    # Bottom of left box — bookings vs revenue
    add_line(slide, 0.70, 6.30, 9.70, 6.30, color=GREY_BORD2, width=0.5)
    add_rich_textbox(slide, 0.70, 6.38, 9.00, 0.45, [
        ("Why this matters for valuation: ", {'size': 9, 'bold': True, 'color': ORANGE_HOT}),
        ("Comp buckets price one or two vectors at a time. No public comp prices all four together — the structural gap that Series G captures.",
         {'size': 9, 'color': TEXT_DARK}),
    ])

    # ============ §2 RIGHT — EXPANSION VECTORS ============
    section_label(slide, 10.10, 2.10, 9.40, "§2  HOW THEY'RE EXPANDING  ·  4 ACTIVE VECTORS WITH DATED EVIDENCE")
    add_rect(slide, 10.10, 2.50, 9.40, 4.40, fill=WHITE, line=GREY_BORD)

    hdr2 = [(0.00, 2.45, "Expansion vector"),
            (2.45, 3.85, "Evidence (dated)"),
            (6.30, 2.90, "What it unlocks")]
    for x_off, w, lbl in hdr2:
        add_textbox(slide, 10.30 + x_off, 2.62, w, 0.25, lbl,
                    size=9.5, bold=True, color=TEXT_GREY)
    add_line(slide, 10.30, 2.92, 19.30, 2.92, color=GREY_BORD2, width=0.5)

    exp = [
        ("Clinical / payer pathway",
         "CMS Innovation Center selection (Apr 13, 2026) · Abbott + Mayo as Series G strategics (Mar 2026) · FDA-cleared ECG · BP Insights launched Jan 2026",
         "Reimbursed-revenue economics; lifts terminal multiple from 3.5× (sub) to 5.0-6.4× (medtech tier)"),
        ("Geographic expansion",
         "~40-45% intl revenue today · 60% of NEW sales intl (NYT DealBook, Mar 2026) · No CAC penalty disclosed",
         "Doubles addressable TAM without product cost; matches Garmin's intl mix"),
        ("Enterprise (B2B Unite)",
         "5% → ~20% of recurring rev over <3 years (Sacra / BMCT) · Per-seat pricing 4-7× consumer",
         "ARPU lift mechanically; channels into payer-paid economics when health-plan contracts open"),
        ("Adjacent biomarkers / AI",
         "Advanced Labs (350K+ waitlist) · BPI launched Jan 2026 · WHOOP Coach (gen AI) · 600+ planned 2026 hires",
         "Premium feature attach lifts ARPU continuously; converts data flywheel into the Bull 1% Labs/BP component"),
    ]
    y = 3.00
    for vec, ev, unlock in exp:
        add_textbox(slide, 10.30,        y, 2.30, 0.90, vec,    size=10, bold=True, color=NAVY)
        add_textbox(slide, 10.30 + 2.45, y, 3.75, 0.90, ev,     size=8.5, color=TEXT_DARK)
        add_textbox(slide, 10.30 + 6.30, y, 3.00, 0.90, unlock, size=8.5, color=TEXT_GREY)
        y += 0.92

    # Bottom of right box — strategic punctuation
    add_line(slide, 10.30, 6.30, 19.30, 6.30, color=GREY_BORD2, width=0.5)
    add_rich_textbox(slide, 10.30, 6.38, 9.00, 0.45, [
        ("Convergence: ", {'size': 9, 'bold': True, 'color': ORANGE_HOT}),
        ("All four expansion vectors feed Bull's ARPU decomposition (Labs / Unite / payer) and Bull's member ceiling (clinical channel TAM). 30% Bull weight = P(all four activate simultaneously).",
         {'size': 9, 'color': TEXT_DARK}),
    ])

    # ============ §3 MIDDLE BAND — THE COMP DISCONNECT QUANTIFIED ============
    section_label(slide, 0.50, 7.10, 19.00, "§3  THE COMP DISCONNECT  ·  WHAT EACH BUCKET MISSES")
    add_rect(slide, 0.50, 7.50, 19.00, 2.10, fill=GREY_BG, line=GREY_BORD2)

    cols3 = [
        ("BUCKET 1 — HARDWARE",       "Garmin · GoPro · Fitbit", "6.2-6.7× rev",
         "Undervalues recurring base · ignores subscription compounding · prices like terminal-growth durables",
         "Implied EV: $557M (T1 median)",  RED),
        ("BUCKET 2 — SUBSCRIPTION",   "Spotify · Roku · Peloton", "1.2-5.1× rev",
         "Doesn't capture healthcare optionality · Peloton drags the floor · Spotify ceiling at 5.1× still leaves gap",
         "Implied EV: $1.74B (T2 median)", NAVY),
        ("BUCKET 3 — HEALTH DATA",    "Dexcom · ResMed · Masimo · iRhythm", "5.0-6.4× rev",
         "Just compressed 50-67% from 2023 peaks · prices reimbursed cash flows that WHOOP doesn't have yet",
         "Implied EV: $2.57B (T3 median)", GREEN),
        ("ORURA — PRIVATE PARALLEL",  "Series E ($12B, Oct 2025)", "12-13× rev",
         "Closest direct comp — but also private · 80% hardware mix vs WHOOP 80% sub · 17% EBITDA margin",
         "Reference only — never blended", NAVY_DEEP),
    ]
    cw = 4.65
    cx = 0.70
    for i, (head, comps, mult, miss, ev, color) in enumerate(cols3):
        x = cx + i*cw
        add_textbox(slide, x, 7.62, cw - 0.10, 0.24, head, size=10, bold=True, color=color)
        add_textbox(slide, x, 7.86, cw - 0.10, 0.20, comps, size=8.5, color=TEXT_GREY)
        add_textbox(slide, x, 8.06, cw - 0.10, 0.24, mult, size=11, bold=True, color=TEXT_DARK)
        add_textbox(slide, x, 8.32, cw - 0.10, 0.80, miss, size=8.5, color=TEXT_DARK)
        add_textbox(slide, x, 9.18, cw - 0.10, 0.24, ev,   size=9, bold=True, color=color)

    # ============ §4 BOTTOM — THE NUMBERED PUNCHLINE ============
    add_rect(slide, 0.50, 9.75, 19.00, 0.85, fill=NAVY_DEEP, line=NAVY_DEEP, no_line=True)
    punch_runs = [
        ("Punchline:  ", {'size': 12, 'bold': True, 'color': ORANGE}),
        ("Series G prices WHOOP at ", {'size': 12, 'color': WHITE}),
        ("~19× recognized rev / ~9× bookings", {'size': 12, 'bold': True, 'color': WHITE}),
        (" — above every public bucket. The premium isn't an error; it is what sophisticated capital pays for the ", {'size': 12, 'color': WHITE}),
        ("four-vector combination", {'size': 12, 'bold': True, 'color': ORANGE}),
        (" no listed company holds together.", {'size': 12, 'color': WHITE}),
    ]
    add_rich_textbox(slide, 0.70, 9.95, 18.60, 0.45, punch_runs)

    return slide


# ===========================================================================
# SLIDE D — VERDICT (replaces Slide 14)
# ===========================================================================
def build_verdict_slide(prs):
    slide = prs.slides.add_slide(BLANK)
    slide_chrome(
        slide,
        top_label="VERDICT",
        title=None,
        title_orange_run=("Series G is fairly priced to ", "mildly underpriced", ".  Hold / accumulate at secondary discount."),
        subtitle="Three lenses say three different things — the spread between them IS the analytical insight. Intrinsic clears Series G by +14%; backward-looking comps drag the blended FF avg to −10%; forward IPO reference class implies +45%.",
        page_label="14 / 17",
        sources="Sources: Internal scenario-weighted DCF · S&P Capital IQ trading multiples · Hiive / Forge secondary order books · Football field methodology · WHOOP CMS Innovation Center press (Apr 13, 2026)"
    )

    # ============ §1 RECONCILIATION TABLE — THE FOUR LENSES ============
    section_label(slide, 0.50, 2.10, 19.00, "§1  RECONCILIATION  ·  WHAT EACH LENS ACTUALLY SAYS")
    add_rect(slide, 0.50, 2.50, 19.00, 2.85, fill=WHITE, line=GREY_BORD)

    # Headers
    rec_hdr = [
        (0.00, 4.30, "Lens"),
        (4.30, 2.10, "EV ($B)"),
        (6.40, 1.90, "vs Series G"),
        (8.30, 10.55, "What this number is telling you"),
    ]
    for x_off, w, lbl in rec_hdr:
        align = PP_ALIGN.LEFT if x_off in (0.00, 8.30) else PP_ALIGN.RIGHT
        add_textbox(slide, 0.70 + x_off, 2.62, w, 0.25, lbl,
                    size=9.5, bold=True, color=TEXT_GREY, align=align)
    add_line(slide, 0.70, 2.92, 19.30, 2.92, color=GREY_BORD2, width=0.5)

    recs = [
        ("Football field weighted avg",       "$9.09B",  "−10.0%",
         "Backward-looking · 35% weight on Bucket-2 comps drags it · cannot price platform-thesis premium",  RED),
        ("Series G market price (Mar 2026)",  "$10.10B", "—",
         "Anchor · sophisticated capital with healthcare strategics (Abbott · Mayo · CMS pathway live)",     TEXT_DARK),
        ("Scenario-weighted intrinsic (HEADLINE)", "$11.56B", "+14.4%",
         "Neutral 20/50/30 conviction · scenario-level weighting captures Jensen Gap · primary recommendation",  GREEN),
        ("Implied IPO 2027 (1.45× median)",   "$14.69B", "+45.4%",
         "Forward · empirical reference-class step-up across 6 consumer/health platform IPOs (24.7% ann. IRR)",  NAVY),
    ]
    y = 3.05
    for label, ev, vs, read, color in recs:
        is_headline = "HEADLINE" in label
        if is_headline:
            add_rect(slide, 0.55, y-0.02, 18.90, 0.50, fill=ORANGE_BG, line=ORANGE_BORD, line_width=0.5)
        add_textbox(slide, 0.70,        y+0.04, 4.30, 0.40, label,
                    size=11, bold=is_headline, color=NAVY if is_headline else TEXT_DARK)
        add_textbox(slide, 0.70 + 4.30, y+0.04, 2.10, 0.40, ev,
                    size=14, bold=True, color=color, align=PP_ALIGN.RIGHT)
        add_textbox(slide, 0.70 + 6.40, y+0.04, 1.90, 0.40, vs,
                    size=12, bold=is_headline, color=color, align=PP_ALIGN.RIGHT)
        add_textbox(slide, 0.70 + 8.30, y+0.04, 10.55, 0.40, read,
                    size=10, color=TEXT_GREY)
        y += 0.56

    # ============ §2 DIRECTIONAL READS (3 columns) ============
    section_label(slide, 0.50, 5.55, 19.00, "§2  THE DIRECTIONAL READ  ·  WHAT TO DO WITH THIS")
    box_top = 5.95
    box_h = 3.20
    bw = 6.20
    gap = 0.20

    # Box 1 — Series G investor
    bx = 0.50
    add_rect(slide, bx, box_top, bw, box_h, fill=GREY_BG, line=NAVY_DEEP)
    add_rect(slide, bx, box_top, bw, 0.40, fill=NAVY_DEEP, line=NAVY_DEEP, no_line=True)
    add_textbox(slide, bx + 0.15, box_top + 0.08, bw - 0.30, 0.25,
                "FOR THE SERIES G INVESTOR",
                size=11, bold=True, color=WHITE)
    add_textbox(slide, bx + 0.15, box_top + 0.50, bw - 0.30, 0.30,
                "Hold / accumulate at secondary discount",
                size=12, bold=True, color=NAVY)
    add_textbox(slide, bx + 0.15, box_top + 0.85, bw - 0.30, 1.30,
                "Intrinsic margin of safety +14.4% over entry mark.  Empirical median 24.7% annualized IRR over 18mo to IPO clears the institutional hurdle.  Hiive order book at $5.44-$6.88 = 39-51% discount to Series G — entry mark is the most expensive price available in the market.",
                size=10, color=TEXT_DARK)
    add_line(slide, bx + 0.15, box_top + 2.30, bx + bw - 0.15, box_top + 2.30, color=GREY_BORD2)
    add_textbox(slide, bx + 0.15, box_top + 2.40, bw - 0.30, 0.25,
                "ACTION", size=9, bold=True, color=ORANGE_HOT)
    add_textbox(slide, bx + 0.15, box_top + 2.65, bw - 0.30, 0.50,
                "Maintain position; add at secondary if liquid.",
                size=11, bold=True, color=NAVY)

    # Box 2 — IPO public investor
    bx = 0.50 + bw + gap
    add_rect(slide, bx, box_top, bw, box_h, fill=GREY_BG, line=NAVY_DEEP)
    add_rect(slide, bx, box_top, bw, 0.40, fill=NAVY_DEEP, line=NAVY_DEEP, no_line=True)
    add_textbox(slide, bx + 0.15, box_top + 0.08, bw - 0.30, 0.25,
                "FOR THE 2027-28 IPO PUBLIC INVESTOR",
                size=11, bold=True, color=WHITE)
    add_textbox(slide, bx + 0.15, box_top + 0.50, bw - 0.30, 0.30,
                "Wait for re-rate signal",
                size=12, bold=True, color=NAVY)
    add_textbox(slide, bx + 0.15, box_top + 0.85, bw - 0.30, 1.30,
                "Reference-class median 1.45× step-up = $14.69B; 75th percentile 1.95× = $19.71B.  Reaching Bucket-3 medtech tier (5.0-6.4× rev) requires post-IPO re-rate driven by reimbursement materialization — realistic but not guaranteed.  CMS Innovation Center selection is the live tracking signal.",
                size=10, color=TEXT_DARK)
    add_line(slide, bx + 0.15, box_top + 2.30, bx + bw - 0.15, box_top + 2.30, color=GREY_BORD2)
    add_textbox(slide, bx + 0.15, box_top + 2.40, bw - 0.30, 0.25,
                "ACTION", size=9, bold=True, color=ORANGE_HOT)
    add_textbox(slide, bx + 0.15, box_top + 2.65, bw - 0.30, 0.50,
                "Underwrite to $14.7B base / $19.7B bull.",
                size=11, bold=True, color=NAVY)

    # Box 3 — What breaks the thesis
    bx = 0.50 + 2*(bw + gap)
    add_rect(slide, bx, box_top, bw, box_h, fill=GREY_BG, line=NAVY_DEEP)
    add_rect(slide, bx, box_top, bw, 0.40, fill=ORANGE_HOT, line=ORANGE_HOT, no_line=True)
    add_textbox(slide, bx + 0.15, box_top + 0.08, bw - 0.30, 0.25,
                "WHAT BREAKS THE THESIS",
                size=11, bold=True, color=WHITE)
    add_textbox(slide, bx + 0.15, box_top + 0.50, bw - 0.30, 0.30,
                "Healthcare optionality stalls",
                size=12, bold=True, color=ORANGE_HOT)
    add_textbox(slide, bx + 0.15, box_top + 0.85, bw - 0.30, 1.30,
                "If CMS pathway closes, Abbott pull-through fails, or FDA blocks BPI follow-on submission → Bull weight collapses → fair value falls to Pessimistic $8.23B (−18.5% vs Series G).  Secondary risk: SoftBank Series F ratchet status unverified — adverse waterfall would compress common-equivalent value below DCF equity.",
                size=10, color=TEXT_DARK)
    add_line(slide, bx + 0.15, box_top + 2.30, bx + bw - 0.15, box_top + 2.30, color=GREY_BORD2)
    add_textbox(slide, bx + 0.15, box_top + 2.40, bw - 0.30, 0.25,
                "TRACK", size=9, bold=True, color=ORANGE_HOT)
    add_textbox(slide, bx + 0.15, box_top + 2.65, bw - 0.30, 0.50,
                "CMS milestones · FDA BPI follow-on · S-1 cap-table disclosure.",
                size=11, bold=True, color=ORANGE_HOT)

    # ============ §3 ROBUSTNESS STRIP (one line) ============
    add_rect(slide, 0.50, 9.40, 19.00, 1.20, fill=NAVY_DEEP, line=NAVY_DEEP, no_line=True)
    add_textbox(slide, 0.70, 9.52, 18.60, 0.28,
                "§3  ROBUSTNESS  ·  CONCLUSION HOLDS ACROSS REASONABLE CONVICTION",
                size=11, bold=True, color=ORANGE)
    # 3 inline mini-stats
    stats = [
        ("PESSIMISTIC  35 / 50 / 15", "$8.23B",  "−18.5%", "Even half Bull weight clears Bear-only by 3.3×"),
        ("NEUTRAL  20 / 50 / 30  (HEADLINE)", "$11.56B", "+14.4%", "Damodaran scenario practice · Oura Series E anchored"),
        ("OPTIMISTIC  15 / 45 / 40", "$13.53B", "+34.0%", "If CMS pathway lifts Bull probability"),
    ]
    sx = 0.70
    sy = 9.85
    sw = 6.20
    for i, (lbl, val, vs, note) in enumerate(stats):
        add_textbox(slide, sx + i*sw, sy,         sw - 0.20, 0.22, lbl,
                    size=9, bold=True, color=ORANGE)
        add_textbox(slide, sx + i*sw, sy + 0.22,  3.20,      0.30, val,
                    size=14, bold=True, color=WHITE)
        add_textbox(slide, sx + i*sw + 3.20, sy + 0.26, sw - 3.40, 0.30, vs,
                    size=11, bold=True, color=ORANGE)
        add_textbox(slide, sx + i*sw, sy + 0.52,  sw - 0.20, 0.22, note,
                    size=8.5, color=WHITE)

    return slide


# ============= BUILD ALL =============
print("Building IPO slide...")
build_ipo_slide(prs)
print("Building Drivers slide...")
build_drivers_slide(prs)
print("Building Business Model slide...")
build_bm_slide(prs)
print("Building Verdict slide...")
build_verdict_slide(prs)

prs.save(OUT)
print(f"\nSaved: {OUT}")
print(f"Total slides now: {len(prs.slides)}")
print(f"New slides appended at positions {len(prs.slides)-3} through {len(prs.slides)}.")
