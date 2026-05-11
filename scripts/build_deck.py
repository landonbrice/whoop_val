"""WHOOP IC deck builder. Reads ALL anchors from audit/model_outputs.json.

17 slides:
  1. Title
  2. Company overview
  3. Methodology preview
  4. Assumption inventory (Section 14)   [NEW]
  5. DCF architecture — Revenue = Members × ARPU; four-lens triangulation
  6. P&L year-by-year + EV→Equity bridge
  7. WACC bottoms-up build
  8. Three scenarios — the dispersion
  9. Scenarios calculator + Jensen's gap  [REPLACED — combined]
 10. Football field
 11. Sensitivity matrix (heatmap + drivers)  [NEW]
 12. Comps positioning (scatter + bracketing)  [REPLACED — visual]
 13. Sensitivities tornado (backup)
 14. Central thesis
 15. Probability weight robustness
 16. Caveats
 17. Q&A defense
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import json

# Colors
NAVY   = RGBColor(0x1F,0x3A,0x5F)
ORANGE = RGBColor(0xE6,0x6E,0x2C)
GRAY   = RGBColor(0x6B,0x72,0x80)
LIGHT  = RGBColor(0xF3,0xF4,0xF6)
DARK   = RGBColor(0x11,0x18,0x27)
RED    = RGBColor(0xC0,0x39,0x2B)
GREEN  = RGBColor(0x2E,0x7D,0x32)
WHITE  = RGBColor(0xFF,0xFF,0xFF)

# Load model outputs
def _intkeys(d):
    if isinstance(d, dict):
        try:
            return {int(k): _intkeys(v) for k,v in d.items()}
        except (ValueError, TypeError):
            return {k: _intkeys(v) for k,v in d.items()}
    if isinstance(d, list):
        return [_intkeys(x) for x in d]
    return d

_M_raw = json.loads(Path("/home/user/whoop_val/audit/model_outputs.json").read_text())
M = {k: _intkeys(v) for k,v in _M_raw.items()}
TOTAL = 17  # slide count

# Convenience accessors
def fmtB(x_m): return f"${x_m/1000:.2f}B"
def fmtBshort(x_m):
    if abs(x_m) >= 1000:
        return f"${x_m/1000:.1f}B"
    return f"${x_m:.0f}M"
def fmtpct(x): return f"{x*100:+.1f}%"
def fmtpctp(x): return f"{x*100:.1f}%"
def fmtM(x_m): return f"${x_m:.0f}M"

OUT = Path("/home/user/whoop_val/Whoop Valuation Master.audited.fast.pptx")
p = Presentation()
p.slide_width  = Inches(20)
p.slide_height = Inches(11.25)
LAYOUT = p.slide_layouts[6]  # blank

V = "/home/user/whoop_val/audit/visuals"

def add_textbox(slide, x, y, w, h, text, size=12, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Calibri"):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = True
    for m in ('margin_left','margin_right','margin_top','margin_bottom'):
        setattr(tf, m, Inches(0))
    p0 = tf.paragraphs[0]
    for idx, line in enumerate(text.split("\n")):
        para = p0 if idx==0 else tf.add_paragraph()
        para.alignment = align
        r = para.add_run(); r.text = line
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
    return tx

def add_rect(slide, x, y, w, h, fill=LIGHT, line=None, line_w=0.5):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh

def add_line(slide, x1, y1, x2, y2, color=GRAY, w=1.0):
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color; ln.line.width = Pt(w)
    return ln

def add_footer(slide, n):
    add_textbox(slide, 0.5, 10.75, 8, 0.25, "BUSN 20410 · Spring 2026 · Yannelis | Landon Brice", size=9, color=GRAY)
    add_textbox(slide, 18.5, 10.75, 1.2, 0.25, f"{n:02d} / {TOTAL}", size=9, color=GRAY, align=PP_ALIGN.RIGHT)

def title_bar(slide, kicker, title, subtitle=None):
    add_textbox(slide, 0.5, 0.35, 19, 0.3, kicker.upper(), size=11, bold=True, color=ORANGE)
    add_textbox(slide, 0.5, 0.7, 19, 0.7, title, size=26, bold=True, color=NAVY)
    if subtitle:
        add_textbox(slide, 0.5, 1.45, 19, 0.4, subtitle, size=13, color=GRAY)
    add_line(slide, 0.5, 1.85, 19.5, 1.85, color=NAVY, w=1.5)

# Pre-compute derived values
sc_neu = M['scen_neu']
sc_pess = M['scen_pess']
sc_opt  = M['scen_opt']
inp_pess = M['input_pess']
inp_neu  = M['input_neu']
inp_opt  = M['input_opt']
prem_intrinsic = sc_neu / M['series_g'] - 1
prem_ff = M['ff_d20']
ff_avg_m = M['ff_d18'] * 1000
gap_neu = sc_neu - inp_neu

# ===========================================================
# SLIDE 1: Title / Hero
# ===========================================================
s = p.slides.add_slide(LAYOUT)
add_rect(s, 0,0,20,11.25, fill=WHITE)
add_textbox(s, 1, 1.5, 18, 0.5, "WHOOP VALUATION", size=14, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_textbox(s, 1, 2.4, 18, 1.5, "Is Series G $10.1B fairly priced?", size=58, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_textbox(s, 1, 4.3, 18, 0.5, "Pressure-testing the March 2026 Series G — five methods, one verdict.", size=18, color=GRAY, align=PP_ALIGN.CENTER)
s.shapes.add_picture(f"{V}/funding_timeline.png", Inches(1.5), Inches(5.1), width=Inches(17))
add_textbox(s, 1, 9.3, 18, 0.45,
            f"Verdict: intrinsic {fmtB(sc_neu)}  |  {fmtpct(prem_intrinsic)} to Series G  |  modestly under-priced",
            size=16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_footer(s, 1)

# ===========================================================
# SLIDE 2: Company overview
# ===========================================================
s = p.slides.add_slide(LAYOUT)
title_bar(s, "Company", "A premium wearable with healthcare ambitions",
          "Where WHOOP sits on the consumer–medical-device spectrum")
cards = [
    ("WHAT", "Hardware-free subscription wearable — band ships with active membership; revenue is recurring, not transactional.\n\n· 2.5M+ paying members (end-2025)\n· $30/month consumer Pro tier\n· FDA-cleared ECG, BPI, Healthspan\n· WHOOP Coach (generative AI)"),
    ("HOW IT GROWS", "103% YoY subscription growth.\nCash-flow positive in 2025.\nFY25 recognized revenue ~$520M.\n\n· ~80% consumer / ~20% B2B (Unite)\n· 40-45% international revenue\n· 60% of NEW members international\n· 5 interim raises 2022-2025"),
    ("WHY HARD TO VALUE", "No public comparable holds together what WHOOP is becoming:\n\n· Hardware (Garmin) — sells boxes\n· Subscription (Spotify/Peloton) — owns content\n· Medical (Dexcom, ResMed) — paid by payers\n\nWHOOP combines all three vectors before the platform is publicly visible. Comps systematically miss this."),
]
for i,(k,b) in enumerate(cards):
    x = 0.5 + i*6.5
    add_rect(s, x, 2.3, 6.2, 5.0, fill=LIGHT, line=GRAY)
    add_textbox(s, x+0.3, 2.45, 5.6, 0.4, k, size=13, bold=True, color=ORANGE)
    add_textbox(s, x+0.3, 2.9, 5.6, 4.3, b, size=12, color=DARK)
add_rect(s, 0.5, 7.6, 19, 2.5, fill=NAVY)
add_textbox(s, 0.8, 7.75, 19, 0.4, "SERIES G — MARCH 31, 2026", size=12, bold=True, color=ORANGE)
g_facts = [("$575M","raised"),("$10.1B","post-money"),("2.8x","step-up vs 2021"),
           ("Collaborative Fund","lead"),("Abbott · Mayo · QIA · Mubadala","strategic + sovereign"),
           ("IPO 2027","explicit trajectory")]
for i,(n,l) in enumerate(g_facts):
    x = 0.8 + i*3.15
    add_textbox(s, x, 8.3, 3.0, 0.7, n, size=22, bold=True, color=WHITE)
    add_textbox(s, x, 9.1, 3.0, 0.5, l, size=11, color=LIGHT)
add_footer(s, 2)

# ===========================================================
# SLIDE 3: Methodology preview
# ===========================================================
s = p.slides.add_slide(LAYOUT)
title_bar(s, "Methodology", "Architecture: scenario-weighted DCF as primary, comps as cross-check",
          "How four valuation lenses triangulate against the Series G mark")
methods = [
    ("INTRINSIC DCF", "Scenario-weighted\n(Bear / Base / Bull)", fmtB(sc_neu), "40% weight"),
    ("PUBLIC COMPS", "Three buckets:\nhardware · sub · health-data", fmtB(M['ff_d10']*1000), "35% weight"),
    ("PRECEDENT M&A", "14 deals across\nthe three buckets", "$5-8B*", "0% (cross-check)"),
    ("LAST-ROUND", "Series G post-money\n(Mar 2026)", "$10.10B", "10% weight"),
    ("IMPLIED IPO", "Private→IPO step-up\nreference class (6 deals)", fmtB(M['ff_d14']*1000), "15% weight"),
]
for i,(k,d,ev,w) in enumerate(methods):
    x = 0.5 + i*3.85
    add_rect(s, x, 2.3, 3.7, 4.0, fill=WHITE, line=NAVY, line_w=1.5)
    add_textbox(s, x+0.2, 2.45, 3.4, 0.4, k, size=11, bold=True, color=ORANGE)
    add_textbox(s, x+0.2, 2.9, 3.4, 1.2, d, size=11, color=DARK)
    add_textbox(s, x+0.2, 4.4, 3.4, 0.6, ev, size=22, bold=True, color=NAVY)
    add_textbox(s, x+0.2, 5.2, 3.4, 0.4, w, size=10, color=GRAY)
add_textbox(s, 0.5, 6.7, 19, 0.4, "WHY THIS ARCHITECTURE", size=12, bold=True, color=ORANGE)
choices = [
    ("Scenario-weighted DCF, not input-level.", "Jensen's inequality: input-averaging at the driver level destroys correlation. WHOOP's drivers co-move within scenarios — scenario-weighted captures the convexity premium worth $" + f"{gap_neu/1000:.1f}B (Slide 9)."),
    ("Real options retained as exhibit, NOT additive.", "Healthcare reimbursement, women's health, glucose monitoring — modeled as correlated outcome clusters (success $5-6B at 15-20% prob; mixed $2-3B at 40-50%; failure $0.5-1B at 30-40%). Cross-validation only."),
    ("Three-bucket comp framework, disaggregated.", "Bucket 2 (consumer subscription): Peloton 1.2x → Spotify 5.1x. Range IS the analysis; averaging destroys it."),
    ("Take-private LBO leg excluded.", "All-VC preferred capital structure offers no debt capacity; sovereign + crossover signal locked-in IPO path; antitrust blocks Apple/Google strategic acquisition."),
    ("Class connections.", "Lecture 1A (DCF mechanics) · 1B (cost of capital, Slide 7) · 3B (real options as exhibit) · 4B (relative valuation)."),
]
y = 7.15
for k,b in choices:
    add_textbox(s, 0.6, y, 4.2, 0.5, k, size=11, bold=True, color=NAVY)
    add_textbox(s, 4.85, y, 14.5, 0.5, b, size=10.5, color=DARK)
    y += 0.55
add_footer(s, 3)

# ===========================================================
# SLIDE 4: Assumption inventory (Section 14) — NEW
# ===========================================================
s = p.slides.add_slide(LAYOUT)
title_bar(s, "Assumption Inventory",
          "20 ProbWeighted drivers — what we sourced, what we stress-tested",
          "Ranked by Bull−Bear spread as % of |Base|; member growth dominates by section")
s.shapes.add_picture(f"{V}/assumption_inventory.png", Inches(0.4), Inches(2.0),
                     width=Inches(19.2))
# Bottom banner — read
inventory = M.get("inventory", [])
top_drivers = sorted(inventory,
                     key=lambda d: abs(((d.get("bull") or 0) - (d.get("bear") or 0)) / max(abs(d.get("base") or 1e-6), 1e-6)),
                     reverse=True)[:5]
top_names = " · ".join((d.get("driver") or "—").split(" — ")[0] for d in top_drivers)
add_rect(s, 0.5, 10.05, 19, 0.65, fill=NAVY)
add_textbox(s, 0.7, 10.15, 19, 0.5,
            f"Read: 6 of the top 10 drivers are member-growth waypoints (MG1-MG8); ARPU growth is rank-13 by design. " +
            f"WACC and Terminal Multiple are the only non-volume drivers in the top 12 — exactly what Members-as-numerator architecture predicts.",
            size=11, color=WHITE)
add_footer(s, 4)

# ===========================================================
# SLIDE 5: DCF architecture — Revenue = Members × ARPU + four-lens
# ===========================================================
s = p.slides.add_slide(LAYOUT)
title_bar(s, "DCF Architecture", "Revenue = Members × ARPU; Members triangulated through four lenses",
          "Heavy analytical effort on Members; controlled discipline on ARPU")
add_rect(s, 0.5, 2.2, 19, 1.0, fill=NAVY)
add_textbox(s, 0.5, 2.45, 19, 0.5, "Revenue  =  Members  ×  ARPU", size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(s, 0.5, 3.5, 9.2, 0.4, "MEMBERS — FOUR-LENS TRIANGULATION (2033 Base ENDPOINT)", size=11, bold=True, color=ORANGE)
add_rect(s, 0.5, 3.9, 9.2, 5.6, fill=LIGHT, line=GRAY)
lenses = [
    ("Lens A — TAM × Terminal Share", "25M premium-WTP TAM × 6%/yr × 30% share", "→ 11.95M"),
    ("Lens B — Peloton stage-shift",   "10.7M base × 1.08 free-hw × 0.97 ARPU drag", "→ 11.21M"),
    ("Lens C — Multi-comp (Apple/Oura)","Three independent comps converge (median)", "→ 10.63M"),
    ("Lens D — Capital-led",           "$1.4B raised → S&M capacity backsolve",      "→ 8.38M"),
    ("Median across four lenses",      "",                                            "→ 10.92M"),
    ("Calibrated Base endpoint (model)","MG-sequence 49/37/25/19/14/11/8/6%",          "→ 10.92M ✓"),
]
y = 4.1
for l,b,o in lenses:
    is_c = "Calibrated" in l or "Median" in l
    add_textbox(s, 0.7, y, 5.0, 0.4, l, size=11, bold=is_c, color=NAVY if is_c else DARK)
    add_textbox(s, 5.7, y, 2.5, 0.4, b, size=10, color=GRAY)
    add_textbox(s, 8.2, y, 1.4, 0.4, o, size=11, bold=True, color=ORANGE if is_c else NAVY)
    y += 0.5
add_textbox(s, 0.7, 8.6, 8.8, 0.6, "All four converge within ±15% of 10.9M Base endpoint.\nThis convergence IS the methodology — not coincidence.",
            size=11, bold=True, color=NAVY)

add_textbox(s, 10.0, 3.5, 9.5, 0.4, "ARPU — ASYMMETRIC DISCIPLINE BY SCENARIO", size=11, bold=True, color=ORANGE)
add_rect(s, 10.0, 3.9, 9.5, 5.6, fill=LIGHT, line=GRAY)
add_textbox(s, 10.2, 4.1, 9.0, 0.4, "2025 anchor: $274 (subscription $258 + Labs $15 + accessories $1)", size=11, color=DARK)
add_textbox(s, 10.2, 4.5, 9.0, 0.4, "Growth rate IS where the scenario differentiates:", size=11, color=DARK)
rows = [("Bear","2.5%/yr","$274 → $334 by 2033","Consumer pricing pressure; Apple Watch undercuts"),
        ("Base","3.0%/yr","$274 → $347","Modest mix shift to higher tiers; Labs at 5-8% attach"),
        ("Bull","6.0%/yr","$274 → $437","Healthcare reimbursement materializes; bundle ramps")]
y = 5.1
for sc,r,traj,basis in rows:
    color = RED if sc=="Bear" else (NAVY if sc=="Base" else GREEN)
    add_textbox(s, 10.2, y, 1.0, 0.4, sc, size=12, bold=True, color=color)
    add_textbox(s, 11.3, y, 1.3, 0.4, r, size=11, bold=True, color=DARK)
    add_textbox(s, 12.7, y, 3.0, 0.4, traj, size=11, color=DARK)
    add_textbox(s, 15.7, y, 3.6, 0.5, basis, size=9, color=GRAY)
    y += 0.65
add_textbox(s, 10.2, 7.3, 9.0, 2.0,
            "WHY ASYMMETRIC?\n\nBear/Base allow ARPU just enough to keep up with inflation + modest tier mix.\n\nBull is where ARPU levers — but ONLY because Bull also requires healthcare reimbursement materialization, which mechanically lifts ARPU through payer-paid pricing. ARPU growth in Bull is a CONSEQUENCE of the scenario, not an independent assumption.",
            size=10.5, color=DARK)
add_footer(s, 5)

# ===========================================================
# SLIDE 6: P&L year-by-year + EV→Equity bridge  [NEW MECHANICS]
# ===========================================================
s = p.slides.add_slide(LAYOUT)
title_bar(s, "DCF Mechanics — Build", "Year-by-year P&L → FCF → discounting → Equity ($M)",
          "Base case; mid-year convention; 8-year explicit period (2026-2033); terminal value at year 8")

# Top: P&L table years 2025-2033
years = M['years']
add_textbox(s, 0.5, 2.1, 19, 0.3, "P&L YEAR-BY-YEAR  (Base case, $M)", size=11, bold=True, color=ORANGE)
# Phase labels
phase_labels = {2025:"Hist", 2026:"Phase 1", 2027:"Phase 1", 2028:"Phase 2", 2029:"Phase 2", 2030:"Phase 2", 2031:"Phase 3", 2032:"Phase 3", 2033:"Phase 3"}
# Table area
tbl_x = 0.5; tbl_y = 2.45; tbl_w = 19; n_yr = len(years)
col_w = (tbl_w - 3.5) / n_yr  # 3.5 in for labels
add_rect(s, tbl_x, tbl_y, tbl_w, 4.3, fill=WHITE, line=GRAY)
# Header row
add_textbox(s, tbl_x+0.1, tbl_y+0.05, 3.4, 0.3, "", size=10, bold=True, color=GRAY)
for i,y in enumerate(years):
    add_textbox(s, tbl_x+3.5+i*col_w, tbl_y+0.05, col_w, 0.3, str(y), size=10, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
# Phase row
add_textbox(s, tbl_x+0.1, tbl_y+0.32, 3.4, 0.25, "Phase", size=8, color=GRAY)
for i,y in enumerate(years):
    add_textbox(s, tbl_x+3.5+i*col_w, tbl_y+0.32, col_w, 0.25, phase_labels[y], size=8, color=GRAY, align=PP_ALIGN.RIGHT)

# Data rows
def make_row(label, vals_dict, fmt_fn, y_pos, bold=False, color=DARK, divider=False):
    add_textbox(s, tbl_x+0.1, y_pos, 3.4, 0.28, label, size=9.5, bold=bold, color=color)
    for i,yr in enumerate(years):
        v = vals_dict[yr] if vals_dict and yr in vals_dict else None
        txt = fmt_fn(v) if v is not None else "—"
        add_textbox(s, tbl_x+3.5+i*col_w, y_pos, col_w, 0.28, txt, size=9.5, bold=bold, color=color, align=PP_ALIGN.RIGHT)

# Helpers
def fM(v): return f"{v:.0f}"
def fpct(v): return f"{v*100:.0f}%"
def fmem(v): return f"{v:.1f}"
def farpu(v): return f"${v:.0f}"

y_pos = tbl_y + 0.65
# Members
make_row("Ending Members (M)", M['rb_end_members'], fmem, y_pos); y_pos += 0.30
# ARPU
make_row("ARPU ($)", M['rb_arpu'], farpu, y_pos); y_pos += 0.30
# Revenue
make_row("Revenue", M['pl_revenue'], fM, y_pos, bold=True, color=NAVY); y_pos += 0.30
# YoY
make_row("  YoY Growth", M['pl_yoy'], fpct, y_pos); y_pos += 0.28
# Gross Profit
make_row("Gross Profit", M['pl_gross'], fM, y_pos); y_pos += 0.28
make_row("  Gross Margin", M['pl_gm'], fpct, y_pos); y_pos += 0.28
# OpEx
make_row("Total OpEx", M['pl_opex'], fM, y_pos); y_pos += 0.28
# EBIT
make_row("EBIT", M['pl_ebit'], fM, y_pos, bold=True); y_pos += 0.28
# Less taxes
make_row("  Less: Cash Taxes", M['fcf_taxes'], fM, y_pos); y_pos += 0.28
# Plus D&A
make_row("  Plus: D&A", M['fcf_da'], fM, y_pos); y_pos += 0.28
# Less CapEx
make_row("  Less: CapEx", M['fcf_capex'], fM, y_pos); y_pos += 0.28
# Less ΔNWC
make_row("  Plus: Δ NWC (def rev)", M['fcf_dnwc'], fM, y_pos); y_pos += 0.28
# UFCF
make_row("Unlevered FCF", M['fcf_ufcf'], fM, y_pos, bold=True, color=ORANGE); y_pos += 0.30
# FCF Margin
make_row("  FCF Margin", M['fcf_margin'], fpct, y_pos); y_pos += 0.28

# Bottom-left: PV table
add_textbox(s, 0.5, 7.0, 11, 0.3, "DISCOUNTING — PV(FCF) AND TERMINAL VALUE", size=11, bold=True, color=ORANGE)
add_rect(s, 0.5, 7.35, 11, 3.1, fill=WHITE, line=GRAY)
# Header
add_textbox(s, 0.7, 7.45, 1.5, 0.25, "Year", size=9.5, bold=True, color=GRAY)
add_textbox(s, 2.3, 7.45, 1.6, 0.25, "UFCF ($M)", size=9.5, bold=True, color=GRAY, align=PP_ALIGN.RIGHT)
add_textbox(s, 4.0, 7.45, 1.5, 0.25, "Period (yr)", size=9.5, bold=True, color=GRAY, align=PP_ALIGN.RIGHT)
add_textbox(s, 5.6, 7.45, 1.7, 0.25, "Discount Factor", size=9.5, bold=True, color=GRAY, align=PP_ALIGN.RIGHT)
add_textbox(s, 7.4, 7.45, 1.7, 0.25, "PV(FCF) ($M)", size=9.5, bold=True, color=GRAY, align=PP_ALIGN.RIGHT)
add_textbox(s, 9.2, 7.45, 2.2, 0.25, f"@ WACC = {M['wacc_base']*100:.2f}%", size=9.5, color=GRAY, align=PP_ALIGN.RIGHT)

forecast_years = [2026,2027,2028,2029,2030,2031,2032,2033]
yp = 7.72
for yr in forecast_years:
    add_textbox(s, 0.7, yp, 1.5, 0.22, str(yr), size=9.5, color=DARK)
    add_textbox(s, 2.3, yp, 1.6, 0.22, f"{M['dcf_ufcf'][yr]:.0f}", size=9.5, color=DARK, align=PP_ALIGN.RIGHT)
    add_textbox(s, 4.0, yp, 1.5, 0.22, f"{M['dcf_discount_pd'][yr]:.1f}", size=9.5, color=DARK, align=PP_ALIGN.RIGHT)
    add_textbox(s, 5.6, yp, 1.7, 0.22, f"{M['dcf_disc_factor'][yr]:.4f}", size=9.5, color=DARK, align=PP_ALIGN.RIGHT)
    add_textbox(s, 7.4, yp, 1.7, 0.22, f"{M['dcf_pv_fcf'][yr]:.0f}", size=9.5, color=DARK, align=PP_ALIGN.RIGHT)
    yp += 0.24

# Sum + TV
add_line(s, 0.7, yp, 11.3, yp, color=GRAY, w=0.8)
yp += 0.05
add_textbox(s, 0.7, yp, 2.0, 0.25, "Σ PV(FCF)", size=10, bold=True, color=DARK)
add_textbox(s, 7.4, yp, 1.7, 0.25, f"{M['dcf_sum_pv_fcf']:.0f}", size=10, bold=True, color=DARK, align=PP_ALIGN.RIGHT)
yp += 0.27
add_textbox(s, 0.7, yp, 2.5, 0.25, f"Terminal Value (yr8)", size=10, color=DARK)
add_textbox(s, 7.4, yp, 1.7, 0.25, f"PV = {M['dcf_pv_tv']:.0f}", size=10, color=DARK, align=PP_ALIGN.RIGHT)
yp += 0.22
add_textbox(s, 0.7, yp, 11, 0.30,
            f"  Exit 3.5×Rev = ${M['dcf_tv_exit']/1000:.1f}B  ·  Gordon (g=3%) = ${M['dcf_tv_gordon']/1000:.1f}B  ·  Blended = ${M['dcf_tv_blend']/1000:.1f}B",
            size=9, color=GRAY)

# Bottom-right: EV→Equity waterfall (drawn as horizontal bar bridge)
add_textbox(s, 12.0, 7.0, 7.5, 0.3, "EV → EQUITY BRIDGE (Base)", size=11, bold=True, color=ORANGE)
add_rect(s, 12.0, 7.35, 7.5, 3.1, fill=WHITE, line=GRAY)

ev = M['dcf_ev']
debt = M['dcf_debt']  # negative number
cash = M['dcf_cash']
equity = M['dcf_equity']

bridge_items = [
    ("Σ PV(FCF)",      M['dcf_sum_pv_fcf'],  NAVY),
    ("+ PV(TV)",       M['dcf_pv_tv'],       NAVY),
    ("= Enterprise V", ev,                   ORANGE),
    ("Less: Debt",     debt,                 RED),
    ("Plus: Cash",     cash,                 GREEN),
    ("= Equity Value", equity,               ORANGE),
]
max_v = max(abs(v) for _,v,_ in bridge_items)
bar_x0 = 14.0
bar_max_w = 5.2
yp2 = 7.55
for lbl, val, color in bridge_items:
    add_textbox(s, 12.1, yp2, 1.8, 0.30, lbl, size=10.5, color=DARK)
    w = bar_max_w * abs(val) / max_v
    add_rect(s, bar_x0, yp2, w, 0.25, fill=color)
    add_textbox(s, bar_x0 + w + 0.1, yp2, 1.4, 0.25, f"{val:.0f}", size=10.5, bold=True, color=color)
    yp2 += 0.42

# Summary callout
add_rect(s, 12.0, 10.45, 7.5, 0.0)  # spacer hidden
add_footer(s, 6)

# ===========================================================
# SLIDE 7: WACC bottoms-up build  [NEW MECHANICS]
# ===========================================================
s = p.slides.add_slide(LAYOUT)
title_bar(s, "DCF Mechanics — WACC", f"Bottoms-up build: {M['wacc_base']*100:.2f}% Base WACC",
          "Comp regressions → unlever → bucket-weight → re-lever → CAPM. Damodaran sector cross-check converges within 0.07 β.")

# Left half: per-comp betas table
add_textbox(s, 0.5, 2.1, 9.2, 0.35, "§1  PER-COMP REGRESSED BETAS  (3Y weekly, Bloomberg-adj)", size=11, bold=True, color=ORANGE)
add_rect(s, 0.5, 2.5, 9.2, 3.2, fill=WHITE, line=GRAY)
add_textbox(s, 0.7, 2.6, 2.5, 0.25, "Comp", size=9.5, bold=True, color=GRAY)
add_textbox(s, 3.3, 2.6, 3.4, 0.25, "Bucket", size=9.5, bold=True, color=GRAY)
add_textbox(s, 6.8, 2.6, 1.2, 0.25, "β raw", size=9.5, bold=True, color=GRAY, align=PP_ALIGN.RIGHT)
add_textbox(s, 8.0, 2.6, 0.7, 0.25, "β adj", size=9.5, bold=True, color=GRAY, align=PP_ALIGN.RIGHT)
add_textbox(s, 8.8, 2.6, 0.7, 0.25, "R²", size=9.5, bold=True, color=GRAY, align=PP_ALIGN.RIGHT)
y = 2.92
for c in M['comps']:
    add_textbox(s, 0.7, y, 2.5, 0.25, c['name'], size=10, color=DARK)
    add_textbox(s, 3.3, y, 3.4, 0.25, c['bucket'], size=10, color=DARK)
    add_textbox(s, 6.8, y, 1.2, 0.25, f"{c['beta_raw']:.3f}", size=10, color=GRAY, align=PP_ALIGN.RIGHT)
    add_textbox(s, 8.0, y, 0.7, 0.25, f"{c['beta_adj']:.3f}", size=10, bold=True, color=DARK, align=PP_ALIGN.RIGHT)
    add_textbox(s, 8.8, y, 0.7, 0.25, f"{c['r2']:.3f}", size=10, color=GRAY, align=PP_ALIGN.RIGHT)
    y += 0.27

add_textbox(s, 0.7, 5.05, 9.0, 0.5, "Notes: All zero-net-debt or near-zero (unlevered = levered). PTON excluded from B2 median (distress).\nR² range 5.7% (IRTC) - 21.6% (GRMN). Bloomberg adjustment: β_adj = ⅔ β_raw + ⅓ × 1.",
            size=8.5, color=GRAY)

# Bucket medians
add_textbox(s, 0.5, 5.85, 9.2, 0.35, "§2  BUCKET MEDIAN UNLEVERED β  (ex-PTON for B2)", size=11, bold=True, color=ORANGE)
add_rect(s, 0.5, 6.25, 9.2, 1.2, fill=WHITE, line=GRAY)
add_textbox(s, 0.7, 6.35, 3.5, 0.25, "Bucket", size=9.5, bold=True, color=GRAY)
add_textbox(s, 4.3, 6.35, 4.0, 0.25, "Comps", size=9.5, bold=True, color=GRAY)
add_textbox(s, 8.4, 6.35, 1.1, 0.25, "Median βᵤ", size=9.5, bold=True, color=GRAY, align=PP_ALIGN.RIGHT)
bucket_rows = [
    ("B1 — Consumer hardware", "Garmin", M['b1_beta']),
    ("B2 — Consumer subscription", "Spotify (ex-Peloton)", M['b2_beta_ex_pton']),
    ("B3 — Health data / med-device", "Dexcom · ResMed · Masimo · iRhythm", M['b3_beta']),
]
y = 6.65
for b,c,beta in bucket_rows:
    add_textbox(s, 0.7, y, 3.5, 0.25, b, size=10, color=DARK)
    add_textbox(s, 4.3, y, 4.0, 0.25, c, size=10, color=GRAY)
    add_textbox(s, 8.4, y, 1.1, 0.25, f"{beta:.3f}", size=10, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
    y += 0.28

# Right half: WACC build-up
add_textbox(s, 10.0, 2.1, 9.5, 0.35, "§3  WHOOP BLENDED β  (re-levered, zero net debt)", size=11, bold=True, color=ORANGE)
add_rect(s, 10.0, 2.5, 9.5, 1.55, fill=WHITE, line=GRAY)
add_textbox(s, 10.2, 2.6, 9.1, 0.25, "Scenario  ·  Bucket weights (B1/B2/B3)  →  Blended βᵤ", size=9.5, bold=True, color=GRAY)
beta_rows = [
    ("Bear",  "40 / 40 / 20", M['beta_blended_bear']),
    ("Base",  "20 / 40 / 40", M['beta_blended_base']),
    ("Bull",  "10 / 30 / 60", M['beta_blended_bull']),
]
y = 2.95
for sc, weights, beta in beta_rows:
    c = RED if sc=="Bear" else (NAVY if sc=="Base" else GREEN)
    add_textbox(s, 10.2, y, 1.0, 0.3, sc, size=11, bold=True, color=c)
    add_textbox(s, 11.3, y, 3.5, 0.3, weights, size=10.5, color=DARK)
    add_textbox(s, 15.5, y, 3.9, 0.3, f"→  βᵤ = {beta:.3f}", size=11, bold=True, color=c)
    y += 0.32

# WACC build
add_textbox(s, 10.0, 4.3, 9.5, 0.35, "§4  CAPM BUILD  ·  Cost of Equity (= WACC, zero debt)", size=11, bold=True, color=ORANGE)
add_rect(s, 10.0, 4.7, 9.5, 4.2, fill=LIGHT, line=NAVY, line_w=1.5)

# Build-up rows for BASE scenario
build_rows = [
    ("Risk-free rate",                   "10Y UST · Fed H.15",                        f"{M['rf']*100:.2f}%",                "+"),
    ("Equity risk premium (rev-wtd)",    "Damodaran Jan 2026 · 42% intl",             f"{M['erp_revwtd_base']*100:.2f}%",   "× β"),
    ("Blended β (Base, ex-PTON)",        "Bucket-weighted from bottom-up",            f"{M['beta_blended_base']:.3f}",      "="),
    ("ERP × β",                          "",                                          f"{M['erp_revwtd_base']*M['beta_blended_base']*100:.2f}%", "+"),
    ("Size premium",                     "Kroll CRSP Decile 3 ($7-13B)",              f"{M['size_prem']*100:.2f}%",         "+"),
    ("Liquidity premium",                "Oura tender (Jan 2026) — 25% DLOM anchor",  f"{M['liq_base']*100:.2f}%",          "="),
]
y = 4.85
for label, source, val, op in build_rows:
    is_total = "ERP × β" in label
    add_textbox(s, 10.2, y, 4.8, 0.3, label, size=10.5, bold=True if is_total else False, color=NAVY if is_total else DARK)
    add_textbox(s, 15.0, y, 3.0, 0.3, source, size=8.5, color=GRAY)
    add_textbox(s, 18.1, y, 0.9, 0.3, val, size=11, bold=True, color=NAVY if is_total else DARK, align=PP_ALIGN.RIGHT)
    add_textbox(s, 19.0, y, 0.4, 0.3, op, size=10, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    y += 0.36

# Final result
add_line(s, 10.2, y+0.05, 19.3, y+0.05, color=NAVY, w=1.5)
y += 0.15
add_textbox(s, 10.2, y, 7.5, 0.35, "WACC (Base) =", size=14, bold=True, color=NAVY)
add_textbox(s, 17.5, y, 1.8, 0.35, f"{M['wacc_base']*100:.2f}%", size=18, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)
y += 0.42
# Scenario WACCs
add_textbox(s, 10.2, y, 9.0, 0.3,
            f"Bear: {M['wacc_bear']*100:.2f}%   ·   Base: {M['wacc_base']*100:.2f}%   ·   Bull: {M['wacc_bull']*100:.2f}%",
            size=10.5, color=DARK)
y += 0.32
add_textbox(s, 10.2, y, 9.0, 0.5, "Bull WACC < Base because Bull weights more into B3 (lower β, longer-duration medtech-style cash flows).", size=9, color=GRAY)

# Cross-check banner at bottom
add_rect(s, 0.5, 9.6, 19, 1.0, fill=NAVY)
add_textbox(s, 0.7, 9.7, 19, 0.3, "CROSS-CHECK · DAMODARAN TOP-DOWN", size=11, bold=True, color=ORANGE)
add_textbox(s, 0.7, 10.0, 19, 0.45,
            f"Damodaran sector approach (4.30% Rf + 0.83 × 4.50% ERP + 0.81% size + 2.00% CSRP) lands at 10.85% — bottoms-up sits 75bps higher due to revenue-weighted intl ERP (+10bps) and beta upgrade vs sector medians (+65bps). Bottoms-up is the conservative choice for the live DCF.",
            size=11, color=WHITE)
add_footer(s, 7)

# ===========================================================
# SLIDE 8: Three scenarios — dispersion
# ===========================================================
s = p.slides.add_slide(LAYOUT)
title_bar(s, "Three Scenarios", f"Bear {fmtBshort(M['bear_equity'])} · Base {fmtBshort(M['base_equity'])} · Bull {fmtBshort(M['bull_equity'])} — feel the dispersion",
          "Each scenario is internally consistent; weights collapse them into a single estimate")
s.shapes.add_picture(f"{V}/scenario_dispersion.png", Inches(0.5), Inches(2.1), width=Inches(11))
add_textbox(s, 11.8, 2.1, 7.8, 0.35, "EACH SCENARIO IS A STORY", size=11, bold=True, color=ORANGE)
stories = [
    ("BEAR — Consumer Wearable Plateau", RED,
     f"Healthcare thesis stalls. Apple Watch + Oura + glucose-sensors split the consumer market. WHOOP becomes Peloton: distressed subscription model, churn rises 12% → 18%, ARPU growth flattens at 2.5%/yr.\n\n2033 ending: 6.5M | Bull-case haircut on every line.\nDCF Equity → {fmtB(M['bear_equity'])}"),
    ("BASE — Durable Consumer Subscription", NAVY,
     f"Healthcare optionality is a real call but doesn't cash. Consumer subscription compounds at moderate rate. Members 2.5M → 10.9M (49%/37%/25%/19% etc.). ARPU 3% growth. Margins ramp to 22% terminal FCF margin.\n\n2033 ending: 10.9M | Mid-case for everything.\nDCF Equity → {fmtB(M['base_equity'])}"),
    ("BULL — Healthcare Reimbursement + Platform Win", GREEN,
     f"FDA expansions + CMS Innovation Center pathway → reimbursement live by 2029. WHOOP Unite enterprise scales. Members 2.5M → 18.1M. ARPU 6%/yr (driven BY healthcare, not assumed independently). Terminal margin 30%.\n\n2033 ending: 18.1M | Bull-stories must align.\nDCF Equity → {fmtB(M['bull_equity'])}"),
]
y = 2.55
for label, color, body in stories:
    add_textbox(s, 11.8, y, 7.8, 0.32, label, size=11, bold=True, color=color)
    add_textbox(s, 11.8, y+0.36, 7.8, 2.3, body, size=10, color=DARK)
    y += 2.65
add_rect(s, 0.5, 9.3, 19, 1.1, fill=LIGHT)
add_textbox(s, 0.7, 9.45, 6, 0.4, "COLLAPSE TO ONE NUMBER", size=11, bold=True, color=ORANGE)
add_textbox(s, 0.7, 9.85, 19, 0.5,
            f"Scenario-weighted Neutral 20/50/30 → {fmtB(sc_neu)} intrinsic. Why scenario-level (not input-level)? Next slide — the novel-grade insight.",
            size=12, color=DARK)
add_footer(s, 8)

# ===========================================================
# SLIDE 9: Scenarios calculator + Jensen's gap (COMBINED)
# ===========================================================
s = p.slides.add_slide(LAYOUT)
title_bar(s, "★ NOVEL · Scenarios & Jensen's Gap",
          f"Three weight schemes; scenario-level beats input-level at every conviction by ${gap_neu/1000:.2f}B+",
          "Same Bear/Base/Bull DCFs. Two ways to combine. The convex curve says the methodology call IS the headline.")

# Top band: Scenarios calculator table — 3 schemes × 6 cols
add_textbox(s, 0.5, 2.05, 19, 0.35,
            "SCENARIOS CALCULATOR — all three weight schemes (live from model)",
            size=11, bold=True, color=ORANGE)
add_rect(s, 0.5, 2.4, 19, 1.95, fill=WHITE, line=GRAY)

ws = M.get("weight_schemes", {})
col_x = [0.7, 5.0, 6.4, 7.8, 10.3, 12.8, 15.2, 17.5]
hdrs = ["Weight scheme", "P(Bear)", "P(Base)", "P(Bull)",
        "Scenario-wtd ($B)", "ProbWtd ($B)", "Jensen gap ($B)", "vs Series G"]
for i, h in enumerate(hdrs):
    add_textbox(s, col_x[i], 2.5, col_x[i+1]-col_x[i] if i < len(col_x)-1 else 2.0,
                0.4, h, size=10, bold=True, color=GRAY)
add_line(s, 0.7, 2.9, 19.3, 2.9, color=GRAY, w=0.8)

scheme_rows = [
    ("pessimistic", "Pessimistic"),
    ("neutral",     "Neutral (HEADLINE)"),
    ("optimistic",  "Optimistic"),
]
y_tbl = 3.0
for key, label in scheme_rows:
    sw = ws.get(key, {})
    is_n = key == "neutral"
    color = ORANGE if is_n else DARK
    cells = [
        label,
        f"{(sw.get('p_bear') or 0)*100:.0f}%",
        f"{(sw.get('p_base') or 0)*100:.0f}%",
        f"{(sw.get('p_bull') or 0)*100:.0f}%",
        f"${(sw.get('scen_wtd') or 0)/1000:.2f}",
        f"${(sw.get('prob_wtd') or 0)/1000:.2f}",
        f"+${(sw.get('jensen') or 0)/1000:.2f}",
        f"{(sw.get('vs_serg') or 0)*100:+.1f}%",
    ]
    for i, v in enumerate(cells):
        w_col = col_x[i+1]-col_x[i] if i < len(col_x)-1 else 2.0
        add_textbox(s, col_x[i], y_tbl, w_col, 0.4, v, size=12,
                    bold=is_n, color=color)
    y_tbl += 0.42

# Main visual: Jensen plot
s.shapes.add_picture(f"{V}/jensens_gap.png", Inches(0.5), Inches(4.5),
                     width=Inches(13))

# Right column: why both, with key takeaway
add_textbox(s, 13.8, 4.5, 5.8, 0.35, "WHY BOTH METHODS",
            size=11, bold=True, color=ORANGE)
add_rect(s, 13.8, 4.9, 5.8, 2.5, fill=LIGHT)
add_textbox(s, 14.0, 5.0, 5.4, 0.35,
            "INPUT-LEVEL (single DCF on prob-wtd inputs)",
            size=10.5, bold=True, color=GRAY)
add_textbox(s, 14.0, 5.4, 5.4, 1.8,
            "Treats drivers as independent.\nValid when scenarios don't co-vary.\nUnderstates convexity when they do.",
            size=10, color=DARK)
add_textbox(s, 14.0, 6.55, 5.4, 0.5,
            f"Headline: {fmtB(inp_neu)}",
            size=14, bold=True, color=GRAY)

add_rect(s, 13.8, 7.55, 5.8, 2.6, fill=LIGHT, line=NAVY, line_w=1.5)
add_textbox(s, 14.0, 7.65, 5.4, 0.35,
            "SCENARIO-LEVEL (P · DCF across scenarios)",
            size=10.5, bold=True, color=NAVY)
add_textbox(s, 14.0, 8.05, 5.4, 1.8,
            "Preserves driver co-movement WITHIN each\nscenario (members↑ → ARPU↑ → margin↑).\nCaptures the convexity premium DCF earns.",
            size=10, color=DARK)
add_textbox(s, 14.0, 9.35, 5.4, 0.5,
            f"Headline: {fmtB(sc_neu)}",
            size=14, bold=True, color=NAVY)

add_rect(s, 0.5, 10.3, 19, 0.45, fill=NAVY)
add_textbox(s, 0.7, 10.36, 19, 0.35,
            f"Headline reading: ${gap_neu/1000:.2f}B Jensen gap at Neutral; +${(ws.get('optimistic',{}).get('jensen') or 0)/1000:.2f}B at Optimistic. " +
            f"The convexity premium IS the mispricing.",
            size=11, bold=True, color=WHITE)
add_footer(s, 9)

# ===========================================================
# SLIDE 10: Football field
# ===========================================================
s = p.slides.add_slide(LAYOUT)
title_bar(s, "Triangulation", "Football field: four methods, two stories",
          "The disconnect IS the analysis — comp methodology pulls weighted avg down; intrinsic captures what comps miss")
s.shapes.add_picture(f"{V}/football_field.png", Inches(0.5), Inches(2.1), width=Inches(13))

add_textbox(s, 13.7, 2.1, 5.8, 0.35, "TWO READS OF THE TRIANGULATION", size=11, bold=True, color=ORANGE)
add_rect(s, 13.7, 2.6, 5.8, 3.5, fill=LIGHT, line=ORANGE, line_w=1.5)
add_textbox(s, 13.9, 2.75, 5.4, 0.35, f"WEIGHTED AVERAGE — {fmtB(ff_avg_m)}", size=12, bold=True, color=ORANGE)
add_textbox(s, 13.9, 3.15, 5.4, 0.35, f"= {fmtB(ff_avg_m)}  →  {fmtpct(prem_ff)} vs Series G", size=11, color=DARK)
add_textbox(s, 13.9, 3.55, 5.4, 2.4,
            f"If you trust the comp pull, Series G is mildly OVERpriced. But:\n\n· Comps weight 35% on a {fmtB(M['ff_d10']*1000)} Bucket-2 (Peloton/Spotify) anchor\n· Comp methodology is structurally backward-looking\n· It cannot price the platform thesis until publicly visible",
            size=10, color=DARK)

add_rect(s, 13.7, 6.3, 5.8, 3.4, fill=LIGHT, line=NAVY, line_w=1.5)
add_textbox(s, 13.9, 6.5, 5.4, 0.35, f"INTRINSIC ALONE — {fmtB(sc_neu)}", size=12, bold=True, color=NAVY)
add_textbox(s, 13.9, 6.9, 5.4, 0.35, f"= {fmtB(sc_neu)}  →  {fmtpct(prem_intrinsic)} vs Series G", size=11, color=DARK)
add_textbox(s, 13.9, 7.3, 5.4, 2.3,
            "Strip the comp drag. The DCF says Series G is meaningfully UNDERpriced.\n\nThe gap between weighted avg and intrinsic IS the structural mispricing of platform-positioned companies — exactly what sophisticated capital captured at Series G entry.\n\nThis is the a16z American Dynamism setup.",
            size=10, color=DARK)
add_footer(s, 10)

# ===========================================================
# SLIDE 11: Sensitivity matrix — weight scheme + material drivers
# ===========================================================
s = p.slides.add_slide(LAYOUT)
title_bar(s, "Sensitivity Matrix",
          "Headline robust to ±15pp on either bear or bull weight; Members drives the rest",
          "Left: scenario-weighted equity across all weight schemes. Right: material drivers (Members · WACC · Terminal multiple).")
s.shapes.add_picture(f"{V}/sensitivity_heatmap.png", Inches(0.4), Inches(2.0),
                     width=Inches(19.2))

# Bottom banner — read
g6 = M.get("grid6", {})
rows = g6.get("rows", [])
# Find range of headline output
if rows:
    flat = [v for r in rows for v in r["equity_b"]]
    g6_min = min(flat)/1000
    g6_max = max(flat)/1000
else:
    g6_min, g6_max = 9.3, 13.8

add_rect(s, 0.5, 10.05, 19, 0.65, fill=NAVY)
add_textbox(s, 0.7, 10.15, 19, 0.5,
            f"Read: scenario-weighted equity ranges ${g6_min:.1f}B – ${g6_max:.1f}B across all 25 weight combinations. "
            f"Even at Pessimistic conviction (P_Bear=30%/P_Bull=20%) the model holds ${rows[-1]['equity_b'][0]/1000:.2f}B — within 8% of Series G. "
            "Headline scenario call is structurally robust to weight-knob variation.",
            size=11, color=WHITE)
add_footer(s, 11)

# ===========================================================
# SLIDE 12: Comps positioning — scatter + bracketing visuals
# ===========================================================
s = p.slides.add_slide(LAYOUT)
title_bar(s, "Comps Disconnect", f"Why comps say {fmtB(M['ff_d10']*1000)} and DCF says {fmtB(sc_neu)}",
          "WHOOP is off the public surface — the absence of a true peer IS the structural mispricing")

# LEFT: scatter (X = NTM growth, Y = EV/Rev)
s.shapes.add_picture(f"{V}/comp_scatter.png", Inches(0.5), Inches(2.05),
                     width=Inches(11.0))

# RIGHT-TOP: bracketing strip
s.shapes.add_picture(f"{V}/comp_bracketing.png", Inches(11.8), Inches(2.05),
                     width=Inches(7.8))

# RIGHT-MID: three-thesis text panel
add_textbox(s, 11.8, 5.7, 7.8, 0.35,
            "THE THREE-BUCKET FRAMEWORK",
            size=11, bold=True, color=ORANGE)
buckets = [
    ("B1 — Consumer hardware",       "Garmin sole; sells boxes."),
    ("B2 — Consumer subscription",   "Peloton (1.2×) ↔ Spotify (3.7×); range IS the insight."),
    ("B3 — Health data / med-device", "Dexcom · iRhythm · ResMed · Masimo — reimbursed; aspirational."),
]
y = 6.1
for k, txt in buckets:
    add_rect(s, 11.8, y, 7.8, 0.85, fill=LIGHT, line=GRAY)
    add_textbox(s, 12.0, y+0.08, 7.4, 0.3, k, size=10.5, bold=True, color=NAVY)
    add_textbox(s, 12.0, y+0.38, 7.4, 0.45, txt, size=9.5, color=DARK)
    y += 0.92

# BOTTOM banner
add_rect(s, 0.5, 8.95, 19, 1.55, fill=NAVY)
add_textbox(s, 0.7, 9.05, 19, 0.4, "THE PLATFORM THESIS",
            size=12, bold=True, color=ORANGE)
add_textbox(s, 0.7, 9.42, 19, 1.05,
            f"WHOOP @ Series G prices at ~19× recognized revenue / ~9× bookings — outside every public tier (highest public anchor: Masimo 5.96×).\n"
            f"No bucket-weighting scheme reaches $10.1B without 70%+ weight on Bucket 3 — which Series G investors implicitly do. " +
            f"WHOOP combines four commercial vectors (consumer sub + clinical-grade hardware + B2B Unite + health-data platform) no public comp holds together. {fmtpct(prem_intrinsic)} intrinsic premium ≈ what gets priced when the fourth vector becomes publicly visible.",
            size=10.5, color=WHITE)
add_footer(s, 12)

# ===========================================================
# SLIDE 13: Sensitivities
# ===========================================================
s = p.slides.add_slide(LAYOUT)
title_bar(s, "Tornado decomposition", "Driver-by-driver: members > terminal multiple > WACC",
          f"Companion to Slide 11 — tornado centered on scenario-weighted Neutral {fmtB(sc_neu)}")
s.shapes.add_picture(f"{V}/tornado.png", Inches(0.5), Inches(2.1), width=Inches(13.5))
add_textbox(s, 14.3, 2.1, 5.2, 0.35, "WHAT THE TORNADO SAYS", size=11, bold=True, color=ORANGE)
reads = [
    ("Members dominates", "2033 ending members (Bear 6.5M ↔ Bull 18.1M) drives ~$8B EV swing. ~3x bigger than next driver. Validates heavy four-lens triangulation focus."),
    ("Terminal multiple #2", "3.0x Bear ↔ 4.5x Bull at $7-8B 2033 revenue → $4-5B swing."),
    ("WACC #3", f"Bear {M['wacc_bear']*100:.2f}% ↔ Bull {M['wacc_bull']*100:.2f}% → ~$3B swing. Bottoms-up build (10.85% Damodaran cross-check) is conservative."),
    ("FCF margin #4", "Terminal 18% Bear ↔ 30% Bull → $2.5B swing. Comparable scale to WACC."),
    ("ARPU growth — surprisingly small", "Rank 5. Confirms ARPU discipline didn't cost the model much."),
    ("Healthcare option prob", "Exhibit-only; 15% materialization shift = $1B EV swing."),
]
y = 2.5
for k,v in reads:
    add_textbox(s, 14.3, y, 5.2, 0.35, k, size=11, bold=True, color=NAVY)
    add_textbox(s, 14.3, y+0.3, 5.2, 0.9, v, size=9.5, color=DARK)
    y += 1.15
add_rect(s, 0.5, 9.6, 19, 0.8, fill=LIGHT)
add_textbox(s, 0.7, 9.75, 19, 0.5,
            f"Joint Members × WACC heatmap (companion exhibit): {fmtB(sc_neu)} Base sits in the dense middle of the surface; even at 20%-pessimistic on Members AND 100bps higher WACC, fair value remains > Series G.",
            size=10.5, color=DARK)
add_footer(s, 13)

# ===========================================================
# SLIDE 14: Central thesis
# ===========================================================
s = p.slides.add_slide(LAYOUT)
title_bar(s, "Recommendation", "Central thesis: Series G is modestly UNDER-priced",
          f"{fmtpct(prem_intrinsic)} intrinsic premium captures structural mispricing of platform-positioned companies")
add_rect(s, 0.5, 2.3, 19, 3.0, fill=NAVY)
add_textbox(s, 0.5, 2.5, 19, 0.5, "INTRINSIC FAIR VALUE", size=14, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_textbox(s, 0.5, 3.0, 19, 1.4, fmtB(sc_neu), size=80, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(s, 0.5, 4.5, 19, 0.5,
            f"vs Series G $10.10B  =  {fmtpct(prem_intrinsic)} premium under Neutral conviction (20/50/30)",
            size=18, color=LIGHT, align=PP_ALIGN.CENTER)
add_textbox(s, 0.5, 5.6, 19, 0.4, "CONVICTION-WEIGHTED ROBUSTNESS", size=12, bold=True, color=ORANGE)
boxes = [
    ("PESSIMISTIC", "35% Bear / 50% Base / 15% Bull", fmtB(sc_pess), fmtpct(sc_pess/M['series_g']-1), "Even if Bull is half as likely as Base assumes, intrinsic stays close to par."),
    ("NEUTRAL (HEADLINE)", "20% Bear / 50% Base / 30% Bull", fmtB(sc_neu), fmtpct(prem_intrinsic), "Damodaran-style scenario weights, anchored against Oura's Series E pricing."),
    ("OPTIMISTIC", "15% Bear / 45% Base / 40% Bull", fmtB(sc_opt), fmtpct(sc_opt/M['series_g']-1), "If healthcare materialization probability rises with CMS Innovation Center pathway."),
]
for i,(k,w,ev,vs,body) in enumerate(boxes):
    x = 0.5 + i*6.5
    is_n = "NEUTRAL" in k
    fill = NAVY if is_n else LIGHT
    tc = WHITE if is_n else DARK
    accent = ORANGE if is_n else NAVY
    add_rect(s, x, 6.0, 6.2, 4.1, fill=fill, line=NAVY, line_w=1.5)
    add_textbox(s, x+0.2, 6.15, 5.8, 0.35, k, size=11, bold=True, color=accent)
    add_textbox(s, x+0.2, 6.55, 5.8, 0.35, w, size=10, color=tc if is_n else GRAY)
    add_textbox(s, x+0.2, 7.0, 5.8, 0.7, ev, size=32, bold=True, color=tc)
    add_textbox(s, x+0.2, 7.85, 5.8, 0.35, f"{vs} vs Series G", size=12, bold=True, color=accent)
    add_textbox(s, x+0.2, 8.4, 5.8, 1.7, body, size=10, color=tc)
add_footer(s, 14)

# ===========================================================
# SLIDE 15: Probability weight robustness
# ===========================================================
s = p.slides.add_slide(LAYOUT)
title_bar(s, "Robustness", "Even at pessimistic conviction, fair value stays close to Series G",
          "Premium is structural, not knob-set; scenario-level dominates input-level at every weight")
add_textbox(s, 0.5, 2.3, 19, 0.4, "WEIGHT SCHEME × METHOD GRID  (intrinsic, $B)", size=11, bold=True, color=ORANGE)
add_rect(s, 0.5, 2.75, 19, 3.8, fill=WHITE, line=GRAY)
hdrs = ["Weight scheme","Bear / Base / Bull","Input-level DCF","Scenario-level DCF","Jensen's gap","Scenario-level vs Series G"]
xs = [0.7, 4.2, 8.5, 11.5, 14.6, 16.8, 19.5]
for i,h in enumerate(hdrs):
    add_textbox(s, xs[i], 2.9, xs[i+1]-xs[i], 0.4, h, size=10, bold=True, color=GRAY)
rows = [
    ("Pessimistic","35 / 50 / 15", inp_pess, sc_pess, sc_pess - inp_pess, sc_pess/M['series_g']-1),
    ("Neutral (headline)","20 / 50 / 30", inp_neu, sc_neu, gap_neu, prem_intrinsic),
    ("Optimistic","15 / 45 / 40", inp_opt, sc_opt, sc_opt - inp_opt, sc_opt/M['series_g']-1),
]
y = 3.5
for sc,wts,inp,scl,gap,vs in rows:
    is_n = "headline" in sc
    color = NAVY if is_n else DARK
    weight = is_n
    cells = [sc, wts, fmtB(inp), fmtB(scl), f"+{fmtB(gap)}", fmtpct(vs)]
    for i,v in enumerate(cells):
        add_textbox(s, xs[i], y, xs[i+1]-xs[i], 0.45, v, size=12, bold=weight, color=color)
    y += 0.7

add_rect(s, 0.5, 6.9, 9, 3.4, fill=LIGHT)
add_textbox(s, 0.7, 7.05, 8.6, 0.4, "READ 1: SCENARIO-LEVEL IS DOMINANT", size=11, bold=True, color=NAVY)
add_textbox(s, 0.7, 7.45, 8.6, 2.7,
            f"Across ALL THREE weight schemes, scenario-level method beats input-level by ${(sc_pess-inp_pess)/1000:.1f}B - ${(sc_opt-inp_opt)/1000:.1f}B. Jensen's gap WIDENS as Bull weight rises.\n\nThis is not a property of the weights. It's a property of the convex DCF. Methodology choice IS the headline.",
            size=11, color=DARK)
add_rect(s, 10.0, 6.9, 9.5, 3.4, fill=LIGHT)
add_textbox(s, 10.2, 7.05, 9.0, 0.4, "READ 2: SERIES G CALL IS ROBUST", size=11, bold=True, color=NAVY)
add_textbox(s, 10.2, 7.45, 9.0, 2.7,
            f"Pessimistic conviction (only 15% Bull weight) → Series G {fmtpct(sc_pess/M['series_g']-1)}.\n\nNeutral (30% Bull) → {fmtpct(prem_intrinsic)} under-priced.\n\nOptimistic (40% Bull) → {fmtpct(sc_opt/M['series_g']-1)} under-priced.\n\nFor Series G to look meaningfully overpriced you must put Bull at ≤10% — Oura at $12B (Series E) and CMS Innovation Center selection make sub-10% Bull hard to defend.",
            size=11, color=DARK)
add_footer(s, 15)

# ===========================================================
# SLIDE 16: Caveats
# ===========================================================
s = p.slides.add_slide(LAYOUT)
title_bar(s, "Caveats", "What the model does NOT capture",
          "Honest scope; the bull case has stipulations")
add_textbox(s, 0.5, 2.3, 9, 0.4, "MODEL LIMITATIONS", size=12, bold=True, color=ORANGE)
caveats_left = [
    ("Healthcare investment cost not modeled.", "Real options as exhibit assume free option; ignore upfront $200-400M clinical/regulatory spend WHOOP must commit."),
    ("Bull case requires explicit story.", "18M ending members presumes healthcare reimbursement live by 2029 — not a residual; an EVENT must occur."),
    ("Series F ratchet status unverified.", "SoftBank absent from Series G syndicate. Ratchet provisions could affect waterfall; could not source disclosure."),
    ("Option pool size opaque.", "Cap IQ shows 11.7% gap between authorized pref and outstanding common; option pool likely 10-15%."),
    ("Cohort churn aged curve directional only.", "WHOOP's 'Pro tier <3% monthly' stat misused industry-wide as blended; pressure-tested via Peloton overlay but no first-principles cohort data."),
]
y = 2.75
for k,v in caveats_left:
    add_textbox(s, 0.7, y, 8.7, 0.4, k, size=11, bold=True, color=NAVY)
    add_textbox(s, 0.7, y+0.3, 8.7, 1.0, v, size=10.5, color=DARK)
    y += 1.35
add_textbox(s, 10.0, 2.3, 9.5, 0.4, "EXTERNAL RISKS NOT IN BASE CASE", size=12, bold=True, color=ORANGE)
caveats_right = [
    ("Apple Watch / Oura / glucose entrants.", "Competitive risk. Bear case captures via 18% terminal churn; could be insufficient if multiple entrants succeed."),
    ("FDA warning letter (Jul 2025) on BPI.", "Healthcare credibility shock; resolved but precedent matters for next FDA submission."),
    ("Whoopgate (May 2025) pricing controversy.", "Estimated $30-40M revenue gap from renewal-timing reversal. Resolved via free upgrades within 48hrs but execution risk persists."),
    ("Founder/key-person risk.", "Will Ahmed concentration. Not directly modeled."),
    ("TV exit-mult vs Gordon divergence (~34%).", "Indicates aggressive 4.5x Bull terminal multiple — stress-tested via tornado, not haircut."),
]
y = 2.75
for k,v in caveats_right:
    add_textbox(s, 10.2, y, 9.0, 0.4, k, size=11, bold=True, color=NAVY)
    add_textbox(s, 10.2, y+0.3, 9.0, 1.0, v, size=10.5, color=DARK)
    y += 1.35
add_rect(s, 0.5, 9.8, 19, 0.65, fill=LIGHT)
add_textbox(s, 0.7, 9.95, 19, 0.5,
            "Real options retained as cross-validation exhibit (correlated cluster: $5-6B success at 15-20% prob). Not added to football field.",
            size=11, color=DARK, align=PP_ALIGN.CENTER)
add_footer(s, 16)

# ===========================================================
# SLIDE 17: Q&A defense
# ===========================================================
s = p.slides.add_slide(LAYOUT)
title_bar(s, "Q&A Defense", "Anticipated questions — backup material",
          "Held in reserve; not necessarily presented unless asked")
qas = [
    ("Defend Bull ARPU 6%/yr — that's well above SaaS norms.",
     "Bull case ARPU growth is a CONSEQUENCE of the scenario, not an independent assumption. Bull requires healthcare reimbursement materialization — payer-paid ARPU mechanically lifts blended ARPU 30-50%. 6%/yr is the geometric average that lands Bull 2033 ARPU at $437 (vs Base $347), consistent with 30-40% mix shift to payer-paid pricing by 2033."),
    ("Why bottoms-up WACC (11.59%) and not Damodaran top-down (10.85%)?",
     "Bottoms-up is the more conservative choice. Both methods converge within ~70bps. Bottoms-up uses our own per-comp regressions (3Y weekly OLS vs S&P 500), Bloomberg-adjusted, with PTON excluded from B2 median due to distress. Top-down Damodaran sector medians (Electronics 0.83, HC Info-Tech 0.99) give 10.85% — slight lower. We use the higher number to keep the headline conservative."),
    ("18M Bull endpoint — defend?",
     "Triangulated by four independent lenses: Lens A (TAM × share) → 11.95M Base / ~16M Bull. Lens B (Peloton stage-shift) → 11.21M Base / ~17M Bull. Lens C (Apple Watch + Oura comp median). Lens D (capital-led: $1.4B raised → ~16M Bull via S&M backsolve). 18M is the upper end of the Bull-aligned envelope, NOT outside it."),
    ("Why 30% Bull weight — not 20%?",
     "Two anchors: (1) Damodaran's standard scenario practice puts upside scenarios at 25-35% when the upside has real-options character. (2) Oura at $12B Series E (Oct 2025) IS market validation that Bull-class wearable + biomarker outcomes are getting funded at >$10B. Discounting Bull to 20% would imply the market is ~50% over-pricing Oura too."),
    ("Why are real options not added to the football field?",
     "Same Jensen logic as the methodology slide. Real options are CORRELATED outcome clusters, not independent additive bets. Healthcare reimbursement, women's health, glucose monitoring co-move within the success/failure clusters. Adding them as a fixed expected value to a separate DCF would double-count the convexity already captured in scenario-weighted Bull. Retained as exhibit; explicitly NOT additive."),
    (f"If Series G is {fmtpct(prem_intrinsic)} under-priced, why didn't more institutional capital pile in?",
     "It DID — 49 investors across 10 rounds; Series G specifically led by Collaborative Fund with Abbott, Mayo, QIA, Mubadala, IVP, Macquarie. Institutional capital captured the platform thesis. The premium is what's left to be re-rated AT IPO once the platform is publicly visible — exactly the a16z American Dynamism arbitrage pattern."),
]
y = 2.3
for k,v in qas:
    add_rect(s, 0.5, y, 19, 1.15, fill=LIGHT, line=GRAY)
    add_textbox(s, 0.7, y+0.05, 18.6, 0.35, k, size=11, bold=True, color=NAVY)
    add_textbox(s, 0.7, y+0.45, 18.6, 0.7, v, size=10, color=DARK)
    y += 1.25
add_footer(s, 17)

p.save(OUT)
print(f"Saved {OUT}: {OUT.stat().st_size:,} bytes, {len(p.slides)} slides")
