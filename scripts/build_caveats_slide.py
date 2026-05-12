"""Caveats slide redesign — categorized by how the model treats each.

3-column architecture replacing the flat 2-column bullet list:
  §1 STRESS-TESTED        — already in tornado/sensitivity
  §2 SCENARIO-ABSORBED    — Bear/Bull stories carry event-dependent risk
  §3 ACKNOWLEDGED         — honest gaps; external / governance
Bottom: real options retained as exhibit (unchanged).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============ TOKENS (matched to existing deck) ============
NAVY        = RGBColor(0x0E, 0x1E, 0x3D)
NAVY_DEEP   = RGBColor(0x1F, 0x3A, 0x5F)
ORANGE      = RGBColor(0xE6, 0x6E, 0x2C)
ORANGE_HOT  = RGBColor(0xC8, 0x44, 0x1A)
ORANGE_NUM  = RGBColor(0xE1, 0x79, 0x3C)
ORANGE_BG   = RGBColor(0xFF, 0xF4, 0xEF)
ORANGE_BORD = RGBColor(0xE8, 0x90, 0x6A)
GREY_BG     = RGBColor(0xF3, 0xF4, 0xF6)
GREY_BG_2   = RGBColor(0xF8, 0xF9, 0xFB)
GREY_BORD   = RGBColor(0x6B, 0x72, 0x80)
GREY_BORD2  = RGBColor(0xE8, 0xE3, 0xDC)
TEXT_DARK   = RGBColor(0x11, 0x18, 0x27)
TEXT_GREY   = RGBColor(0x6B, 0x72, 0x80)
TEXT_GREY2  = RGBColor(0x70, 0x70, 0x70)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NAVY_BG     = RGBColor(0xEC, 0xEF, 0xF4)   # very light navy tint
NAVY_BORD   = RGBColor(0xB8, 0xC2, 0xD4)   # navy chip border

FONT = "Calibri"
FONT_DISPLAY = "Georgia"

# ============ HELPERS ============
def add_textbox(slide, left, top, width, height, text, *,
                font=FONT, size=10, bold=False, italic=False, color=TEXT_DARK,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02);  tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.name = font; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_rich_textbox(slide, left, top, width, height, runs, *,
                     align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02);  tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    for txt, attrs in runs:
        run = p.add_run(); run.text = txt
        f = run.font
        f.name = attrs.get('font', FONT)
        f.size = Pt(attrs.get('size', 10))
        f.bold = attrs.get('bold', False)
        f.italic = attrs.get('italic', False)
        if 'color' in attrs:
            f.color.rgb = attrs['color']
    return tb

def add_rect(slide, left, top, width, height, *,
             fill=WHITE, line=GREY_BORD, line_width=0.5, no_line=False,
             shape=MSO_SHAPE.RECTANGLE):
    rect = slide.shapes.add_shape(shape,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    if fill is None:
        rect.fill.background()
    else:
        rect.fill.solid(); rect.fill.fore_color.rgb = fill
    if no_line:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line; rect.line.width = Pt(line_width)
    rect.shadow.inherit = False
    rect.text_frame.text = ""
    return rect

def add_line(slide, x1, y1, x2, y2, *, color=NAVY_DEEP, width=0.75):
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1),
                                     Inches(x2), Inches(y2))
    ln.line.color.rgb = color; ln.line.width = Pt(width)
    return ln

# ============ BUILD ============
prs = Presentation()
prs.slide_width = Inches(20.0)
prs.slide_height = Inches(11.25)
BLANK = prs.slide_layouts[6]
slide = prs.slides.add_slide(BLANK)

# ---- TOP CHROME (matches existing slide 15) ----
add_textbox(slide, 0.50, 0.35, 19.00, 0.30,
            "CAVEATS", size=11, bold=True, color=ORANGE)

# Title preserved
add_textbox(slide, 0.50, 0.70, 19.00, 0.65,
            "What the model does NOT capture.",
            font=FONT_DISPLAY, size=30, bold=True, color=NAVY)

# Subtitle — REFRAMED: most caveats are already captured somewhere
add_textbox(slide, 0.50, 1.45, 19.00, 0.40,
            "Three honest categories. Most named caveats already sit inside the sensitivity grid or are absorbed by scenario weighting — the genuine gaps are the external risks no DCF can price.",
            size=13, color=TEXT_GREY)

# Divider
add_line(slide, 0.50, 1.95, 19.50, 1.95, color=NAVY_DEEP, width=0.75)

# Footer
add_textbox(slide, 0.50, 10.78, 17.80, 0.22,
            "BUSN 20410  ·  Spring 2026  ·  Yannelis | Landon Brice",
            size=9, color=TEXT_GREY2)
add_textbox(slide, 18.50, 10.75, 1.20, 0.25,
            "16 / 17", size=9, color=TEXT_GREY, align=PP_ALIGN.RIGHT)

# ===========================================================================
# THREE COLUMNS
# ===========================================================================
COL_TOP = 2.20
COL_H = 7.70
COL_W = 6.20
COL_GAP = 0.20

col_x = [0.50, 0.50 + COL_W + COL_GAP, 0.50 + 2*(COL_W + COL_GAP)]
# verify: 0.50 + 3*6.20 + 2*0.20 = 0.50 + 18.60 + 0.40 = 19.50 ✓

# ---- COLUMN 1: STRESS-TESTED (orange tag) ----
x = col_x[0]
# Tag bar
add_rect(slide, x, COL_TOP, COL_W, 0.50, fill=ORANGE_BG, line=ORANGE_BORD,
         line_width=0.75)
# Tag pill (small orange chip)
add_rect(slide, x + 0.20, COL_TOP + 0.13, 1.30, 0.26, fill=ORANGE_HOT,
         line=ORANGE_HOT, no_line=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, x + 0.20, COL_TOP + 0.13, 1.30, 0.26,
            "STRESS-TESTED", size=8.5, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_textbox(slide, x + 1.60, COL_TOP + 0.13, COL_W - 1.80, 0.26,
            "Already in the tornado / sensitivity grid",
            size=10, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

# Body container
add_rect(slide, x, COL_TOP + 0.50, COL_W, COL_H - 0.50,
         fill=WHITE, line=GREY_BORD2, line_width=0.75)

# Caveats in column 1
caveats_1 = [
    ("Member endpoint dominates the answer.",
     "Bear 6.0M ↔ Bull 16.6M swings DCF Equity by ±$3.7B — rank-1 driver, 3× next-largest.",
     "±$3.7B"),
    ("Terminal multiple — Bull 4.5×.",
     "TV exit-multiple vs Gordon-growth divergence 34% (above 25% threshold). Stress-tested via tornado, not haircut.",
     "±$1.5B"),
    ("WACC — bottoms-up vs top-down.",
     "Internal 11.59% sits 75bps above Damodaran top-down 10.85%. Conservative anchor; spread quantified.",
     "±$1.5B"),
    ("ARPU Bull 6% decomposition.",
     "0.7% payer-paid uplift is the loosest of the 5 components — plausible 0.5-1.0% range.",
     "±$0.8B"),
]

card_y = COL_TOP + 0.65
card_h = 1.65
for head, body, mag in caveats_1:
    # Subtle inner card
    add_rect(slide, x + 0.15, card_y, COL_W - 0.30, card_h,
             fill=WHITE, line=ORANGE_BORD, line_width=0.5)
    # Left accent bar
    add_rect(slide, x + 0.15, card_y, 0.08, card_h,
             fill=ORANGE_HOT, line=ORANGE_HOT, no_line=True)
    # Magnitude badge top-right
    add_rect(slide, x + COL_W - 1.20, card_y + 0.12, 0.95, 0.28,
             fill=ORANGE_BG, line=ORANGE_BORD, line_width=0.5,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, x + COL_W - 1.20, card_y + 0.12, 0.95, 0.28,
                mag, size=9.5, bold=True, color=ORANGE_HOT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Headline
    add_textbox(slide, x + 0.35, card_y + 0.12, COL_W - 1.65, 0.32,
                head, size=11, bold=True, color=NAVY)
    # Body
    add_textbox(slide, x + 0.35, card_y + 0.50, COL_W - 0.55, 1.05,
                body, size=9.5, color=TEXT_GREY)
    card_y += card_h + 0.10

# ---- COLUMN 2: SCENARIO-ABSORBED (navy tag) ----
x = col_x[1]
add_rect(slide, x, COL_TOP, COL_W, 0.50, fill=NAVY_BG, line=NAVY_BORD,
         line_width=0.75)
add_rect(slide, x + 0.20, COL_TOP + 0.13, 1.55, 0.26, fill=NAVY_DEEP,
         line=NAVY_DEEP, no_line=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, x + 0.20, COL_TOP + 0.13, 1.55, 0.26,
            "SCENARIO-ABSORBED", size=8.5, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_textbox(slide, x + 1.85, COL_TOP + 0.13, COL_W - 2.05, 0.26,
            "Bear / Bull stories carry the event risk",
            size=10, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

add_rect(slide, x, COL_TOP + 0.50, COL_W, COL_H - 0.50,
         fill=WHITE, line=GREY_BORD2, line_width=0.75)

caveats_2 = [
    ("Bull 18M members presumes reimbursement live by 2029.",
     "Not a residual — an EVENT must occur (FDA + CMS pathway). Priced at 30% Bull probability, anchored externally vs 9% reference-class base rate.",
     "P = 30%"),
    ("Healthcare regulatory / clinical spend $200-400M.",
     "Implicit in Bull's operating thesis but not line-itemed in DCF. Treated as the 'cost of buying the option' — defensible but Yannelis may press.",
     "Bull-only"),
    ("Probability weight selection IS an analytical choice.",
     "Pessimistic 35/50/15 → $8.23B (−18.5%) · Neutral 20/50/30 → $11.56B (+14.4%) · Optimistic 15/45/40 → $13.53B (+34%). Conclusion holds across.",
     "Range shown"),
    ("Healthcare timing uncertainty (post-2033 tail).",
     "Bull case prices materialization within forecast horizon. Delays beyond 2033 reduce DCF contribution — absorbed by Bull probability, not separately discounted.",
     "Long-tail"),
]

card_y = COL_TOP + 0.65
for head, body, mag in caveats_2:
    add_rect(slide, x + 0.15, card_y, COL_W - 0.30, card_h,
             fill=WHITE, line=NAVY_BORD, line_width=0.5)
    add_rect(slide, x + 0.15, card_y, 0.08, card_h,
             fill=NAVY_DEEP, line=NAVY_DEEP, no_line=True)
    add_rect(slide, x + COL_W - 1.20, card_y + 0.12, 0.95, 0.28,
             fill=NAVY_BG, line=NAVY_BORD, line_width=0.5,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, x + COL_W - 1.20, card_y + 0.12, 0.95, 0.28,
                mag, size=9.5, bold=True, color=NAVY_DEEP,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, x + 0.35, card_y + 0.12, COL_W - 1.65, 0.32,
                head, size=11, bold=True, color=NAVY)
    add_textbox(slide, x + 0.35, card_y + 0.50, COL_W - 0.55, 1.05,
                body, size=9.5, color=TEXT_GREY)
    card_y += card_h + 0.10

# ---- COLUMN 3: ACKNOWLEDGED (gray tag — honest gaps) ----
x = col_x[2]
add_rect(slide, x, COL_TOP, COL_W, 0.50, fill=GREY_BG, line=GREY_BORD2,
         line_width=0.75)
add_rect(slide, x + 0.20, COL_TOP + 0.13, 1.45, 0.26, fill=GREY_BORD,
         line=GREY_BORD, no_line=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
add_textbox(slide, x + 0.20, COL_TOP + 0.13, 1.45, 0.26,
            "ACKNOWLEDGED", size=8.5, bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_textbox(slide, x + 1.75, COL_TOP + 0.13, COL_W - 1.95, 0.26,
            "Honest gaps — no DCF can price these",
            size=10, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

add_rect(slide, x, COL_TOP + 0.50, COL_W, COL_H - 0.50,
         fill=WHITE, line=GREY_BORD2, line_width=0.75)

caveats_3 = [
    ("Competitive entrants (Apple / Oura / glucose).",
     "Bear case 18% terminal churn is a partial offset. Could be insufficient if multiple entrants succeed simultaneously — most external risk.",
     "Bear partial"),
    ("Cap structure unknowns.",
     "Series F ratchet status unverified (SoftBank absent from Series G). Option pool likely 10-15% per Cap IQ gap. Adverse waterfall could compress common-equivalent.",
     "Pre-S-1"),
    ("Founder + execution risk.",
     "Will Ahmed concentration not directly modeled. Whoopgate (May 2025) and FDA BPI warning (Jul 2025) precedent operational risk that current architecture can't price.",
     "Qualitative"),
    ("Cohort churn directional only.",
     "WHOOP's 'Pro tier <3% monthly' stat misused industry-wide as blended; pressure-tested via Peloton overlay but no first-principles cohort data until S-1.",
     "Data-gap"),
]

card_y = COL_TOP + 0.65
for head, body, mag in caveats_3:
    add_rect(slide, x + 0.15, card_y, COL_W - 0.30, card_h,
             fill=WHITE, line=GREY_BORD2, line_width=0.5)
    add_rect(slide, x + 0.15, card_y, 0.08, card_h,
             fill=GREY_BORD, line=GREY_BORD, no_line=True)
    add_rect(slide, x + COL_W - 1.20, card_y + 0.12, 0.95, 0.28,
             fill=GREY_BG, line=GREY_BORD2, line_width=0.5,
             shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_textbox(slide, x + COL_W - 1.20, card_y + 0.12, 0.95, 0.28,
                mag, size=9.5, bold=True, color=TEXT_GREY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, x + 0.35, card_y + 0.12, COL_W - 1.65, 0.32,
                head, size=11, bold=True, color=NAVY)
    add_textbox(slide, x + 0.35, card_y + 0.50, COL_W - 0.55, 1.05,
                body, size=9.5, color=TEXT_GREY)
    card_y += card_h + 0.10

# ===========================================================================
# BOTTOM STRIP — REAL OPTIONS RETAINED (preserved from current slide)
# ===========================================================================
strip_top = 9.95
strip_h = 0.60
add_rect(slide, 0.50, strip_top, 19.00, strip_h,
         fill=GREY_BG, line=GREY_BORD2, line_width=0.5)

bottom_runs = [
    ("Real options retained as cross-validation exhibit  ·  ",
     {'size': 11, 'bold': True, 'color': NAVY}),
    ("correlated cluster: $5-6B success at 15-20% probability  ·  ",
     {'size': 11, 'color': TEXT_DARK}),
    ("NOT added to football field (double-counts with Bull bucket-3 multiple weighting).",
     {'size': 11, 'italic': True, 'color': TEXT_GREY}),
]
add_rich_textbox(slide, 0.70, strip_top + 0.15, 18.60, 0.32, bottom_runs,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ============ SAVE ============
OUT = "/home/user/whoop_val/Whoop Pitch - Caveats redesign.pptx"
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)} · Shapes on slide 1: {len(prs.slides[0].shapes)}")
