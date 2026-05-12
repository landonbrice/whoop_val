"""Build a standalone PPTX for the comps slide with the Option-A triangulation strip.

Produces two slides:
  1. Full slide mock — headline + chart placeholder + platform-thesis sidebar + triangulation strip
  2. Triangulation strip alone — drop-in element to paste over the existing deck slide
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# Palette matched to existing deck
NAVY   = RGBColor(0x1F, 0x3A, 0x5F)
ORANGE = RGBColor(0xE6, 0x6E, 0x2C)
GRAY   = RGBColor(0x6B, 0x72, 0x80)
LIGHT  = RGBColor(0xF3, 0xF4, 0xF6)
DARK   = RGBColor(0x11, 0x18, 0x27)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
RULE   = RGBColor(0xD1, 0xD5, 0xDB)

OUT = Path("/home/user/whoop_val/output/whoop-comp-triangulation-slide.pptx")

p = Presentation()
p.slide_width  = Inches(20)
p.slide_height = Inches(11.25)
LAYOUT = p.slide_layouts[6]  # blank


def add_textbox(slide, x, y, w, h, text, size=12, bold=False, color=DARK,
                align=PP_ALIGN.LEFT, font="Calibri", anchor=None, italic=False):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    for m in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, m, Inches(0))
    if anchor is not None:
        tf.vertical_anchor = anchor
    p0 = tf.paragraphs[0]
    for idx, line in enumerate(text.split("\n")):
        para = p0 if idx == 0 else tf.add_paragraph()
        para.alignment = align
        r = para.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tx


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    return shp


def add_line(slide, x1, y1, x2, y2, color=GRAY, weight=1.25, dash=False):
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    if dash:
        try:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            ln.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        except Exception:
            pass
    return ln


def add_oval(slide, cx, cy, d, fill=ORANGE, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(cx - d / 2), Inches(cy - d / 2),
                                 Inches(d), Inches(d))
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    return shp


# =================================================================
# TRIANGULATION STRIP — reusable builder
# =================================================================
def build_strip(slide, x0, y0, w, h):
    """Render the comp-implied valuation strip inside the box (x0,y0,w,h)."""
    # Container card
    add_rect(slide, x0, y0, w, h, fill=LIGHT, line=RULE)

    # Eyebrow / header
    add_textbox(slide, x0 + 0.25, y0 + 0.18, w - 0.5, 0.32,
                "COMP-IMPLIED VALUATION RANGE",
                size=12, bold=True, color=ORANGE)
    add_textbox(slide, x0 + 0.25, y0 + 0.50, w - 0.5, 0.34,
                "Static comps converge $5.1–5.9B. Series G prices $10.1B. The ~$4.2B gap is what the DCF + real-options must defend.",
                size=11.5, color=NAVY, italic=True)

    # Axis geometry
    ax_y     = y0 + h - 0.85          # baseline
    ax_left  = x0 + 0.95
    ax_right = x0 + w - 0.55
    ax_w     = ax_right - ax_left
    vmin, vmax = 0.0, 12.0            # $B
    def vx(v): return ax_left + (v - vmin) / (vmax - vmin) * ax_w

    # Axis line
    add_line(slide, ax_left, ax_y, ax_right, ax_y, color=GRAY, weight=1.5)

    # Tick labels (every $2B)
    for tv in [0, 2, 4, 6, 8, 10, 12]:
        tx = vx(tv)
        add_line(slide, tx, ax_y, tx, ax_y + 0.10, color=GRAY, weight=0.75)
        add_textbox(slide, tx - 0.45, ax_y + 0.13, 0.9, 0.25,
                    f"${tv}B", size=9, color=GRAY, align=PP_ALIGN.CENTER)

    # Static-comp band shading ($5.1B → $5.9B)
    band_x1 = vx(5.1)
    band_x2 = vx(5.9)
    add_rect(slide, band_x1, ax_y - 0.55, band_x2 - band_x1, 0.55,
             fill=RGBColor(0xFF, 0xE9, 0xD6), line=None)
    add_textbox(slide, band_x1 - 0.85, ax_y - 0.78, band_x2 - band_x1 + 1.7, 0.22,
                "STATIC-COMP BAND", size=8, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    # Anchors — (value, label_line1, label_line2, color, above?, weight_str)
    anchors_above = [
        # Above the axis
        (5.1, "Bear",       "40/40/20  ·  4.6×",  NAVY,   True),
        (5.6, "Base",       "20/40/40  ·  5.1×",  NAVY,   True),
        (5.9, "Bull",       "10/30/60  ·  5.4×",  NAVY,   True),
    ]
    anchors_below = [
        # Below the axis — range bars + Series G
        ("range",  5.0, 8.0, "Precedent M&A (14 deals, current-regime medtech 6.6–10.1×)",  GRAY),
        ("range",  5.0, 7.0, "Secondary market — Hiive order book (39–51% discount)",        GRAY),
        ("point", 10.1, None, "Series G  ·  $10.1B  ·  9.2× bookings / 19.4× recognized",   ORANGE),
    ]

    # Above-axis points
    for v, l1, l2, color, _ in anchors_above:
        cx = vx(v)
        add_line(slide, cx, ax_y - 0.55, cx, ax_y, color=color, weight=1.25)
        add_oval(slide, cx, ax_y, 0.16, fill=color)
        add_textbox(slide, cx - 1.0, ax_y - 1.05, 2.0, 0.26,
                    f"${v:.1f}B", size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(slide, cx - 1.0, ax_y - 0.83, 2.0, 0.22,
                    l1, size=9.5, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        # l2 sits a touch above the dot for sub-label
        # (kept off to avoid crowding)

    # Sub-labels for the three bucket points (compact, below the band)
    # Re-render below the axis on a stacked legend row
    legend_y = ax_y + 0.50
    legend_items = [
        ("Bear  40/40/20",  "blended 4.6×",  NAVY),
        ("Base  20/40/40",  "blended 5.1×",  NAVY),
        ("Bull  10/30/60",  "blended 5.4×",  NAVY),
    ]
    leg_x = x0 + 0.30
    leg_w = 2.45
    for i, (k, v, c) in enumerate(legend_items):
        lx = leg_x + i * (leg_w + 0.05)
        add_oval(slide, lx + 0.12, legend_y + 0.10, 0.14, fill=c)
        add_textbox(slide, lx + 0.30, legend_y, 2.2, 0.22, k,
                    size=9, bold=True, color=DARK)
        add_textbox(slide, lx + 0.30, legend_y + 0.20, 2.2, 0.22, v,
                    size=8.5, color=GRAY)

    # Below-axis: range bars
    bar_offsets = {0: 0.95, 1: 1.50}  # vertical offsets below axis for the two ranges
    for idx, item in enumerate([a for a in anchors_below if a[0] == "range"]):
        _, v1, v2, lbl, color = item
        x1 = vx(v1); x2 = vx(v2)
        yb = ax_y + bar_offsets[idx]
        # bar
        add_rect(slide, x1, yb, x2 - x1, 0.10, fill=color, line=None)
        # end caps
        add_line(slide, x1, yb - 0.05, x1, yb + 0.15, color=color, weight=1)
        add_line(slide, x2, yb - 0.05, x2, yb + 0.15, color=color, weight=1)
        # endpoint labels
        add_textbox(slide, x1 - 0.45, yb - 0.02, 0.9, 0.22,
                    f"${v1:.0f}B", size=8.5, color=color, align=PP_ALIGN.CENTER)
        add_textbox(slide, x2 - 0.45, yb - 0.02, 0.9, 0.22,
                    f"${v2:.0f}B", size=8.5, color=color, align=PP_ALIGN.CENTER)
        # label to the right
        add_textbox(slide, x2 + 0.20, yb - 0.04, 6.5, 0.26, lbl,
                    size=9.5, color=DARK)

    # Series G point (orange, anchored on axis with vertical drop)
    for item in anchors_below:
        if item[0] != "point":
            continue
        _, v, _, lbl, color = item
        cx = vx(v)
        # drop-line from axis up to a marker callout
        add_line(slide, cx, ax_y - 0.55, cx, ax_y, color=color, weight=1.5)
        add_oval(slide, cx, ax_y, 0.22, fill=color)
        add_textbox(slide, cx - 1.0, ax_y - 1.05, 2.0, 0.26,
                    f"${v:.1f}B", size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(slide, cx - 1.4, ax_y - 0.83, 2.8, 0.22,
                    "Series G mark", size=10, bold=True, color=color, align=PP_ALIGN.CENTER)

    # GAP bracket — from $5.9B (top of static band) to $10.1B
    gap_x1 = vx(5.9)
    gap_x2 = vx(10.1)
    gap_y  = ax_y - 1.55
    add_line(slide, gap_x1, gap_y, gap_x2, gap_y, color=ORANGE, weight=1.5, dash=True)
    add_line(slide, gap_x1, gap_y, gap_x1, gap_y + 0.18, color=ORANGE, weight=1.5)
    add_line(slide, gap_x2, gap_y, gap_x2, gap_y + 0.18, color=ORANGE, weight=1.5)
    mid_x = (gap_x1 + gap_x2) / 2
    add_textbox(slide, mid_x - 2.5, gap_y - 0.42, 5.0, 0.30,
                "GAP  ≈  $4.2B",
                size=12, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_textbox(slide, mid_x - 4.5, gap_y - 0.15, 9.0, 0.30,
                "bridged by ~1.7× growth premium (103% vs. comps 10–26%) + real-option layer",
                size=10, italic=True, color=DARK, align=PP_ALIGN.CENTER)


# =================================================================
# SLIDE 1: FULL SLIDE MOCK (chart placeholder + sidebar + strip)
# =================================================================
s = p.slides.add_slide(LAYOUT)

# Eyebrow
add_textbox(s, 0.45, 0.30, 18, 0.32,
            "COMPARABLE COMPANIES DISCONNECT",
            size=12, bold=True, color=ORANGE)
# Headline (two lines)
add_textbox(s, 0.45, 0.62, 18.5, 1.10,
            "Series G prices outside every public-comp thesis — top of healthcare is 5.96×,",
            size=26, bold=True, color=NAVY)
add_textbox(s, 0.45, 1.02, 18.5, 1.10,
            "WHOOP bookings is 9.2× pricing near OURA.",
            size=26, bold=True, color=NAVY)
# Rule under headline
add_line(s, 0.45, 1.78, 19.55, 1.78, color=NAVY, weight=1.5)

# Chart placeholder (the existing scatter goes here)
CHART_X, CHART_Y, CHART_W, CHART_H = 0.45, 2.00, 13.6, 5.55
add_rect(s, CHART_X, CHART_Y, CHART_W, CHART_H, fill=WHITE, line=RULE)
add_textbox(s, CHART_X, CHART_Y + CHART_H/2 - 0.30, CHART_W, 0.32,
            "[ paste existing scatter here  ·  EV / NTM Revenue × NTM Growth ]",
            size=12, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, CHART_X, CHART_Y + CHART_H/2 + 0.05, CHART_W, 0.28,
            "tier-coded dots, WHOOP @ 9.2× / 19.4× callout, OURA anchor at 120% growth",
            size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

# Platform-thesis sidebar (kept from existing slide; reproduced here for context)
SB_X, SB_Y, SB_W, SB_H = 14.20, 2.00, 5.35, 5.55
add_rect(s, SB_X, SB_Y, SB_W, SB_H, fill=NAVY)
add_textbox(s, SB_X + 0.25, SB_Y + 0.20, SB_W - 0.5, 0.32,
            "THE PLATFORM THESIS", size=12, bold=True, color=ORANGE)
sidebar_body = (
    "WHOOP @ Series G prices at ~19× recognized revenue / ~9× bookings — outside every "
    "public tier (highest public anchor: Masimo 5.96×).\n\n"
    "No bucket-weighting scheme reaches $10.1B without 70%+ weight on Bucket 3 — which "
    "Series G investors implicitly do.\n\n"
    "WHOOP combines four commercial vectors (consumer sub + clinical-grade hardware + "
    "B2B Unite + health-data platform) no public comp holds together. +14.4% intrinsic "
    "premium ≈ what gets priced when the fourth vector becomes publicly visible."
)
add_textbox(s, SB_X + 0.25, SB_Y + 0.62, SB_W - 0.5, SB_H - 0.8,
            sidebar_body, size=10.5, color=WHITE)

# === NEW: Triangulation strip ===
build_strip(s, x0=0.45, y0=7.70, w=19.10, h=3.30)

# Footer
add_textbox(s, 0.45, 11.00, 19, 0.22,
            "Sources: S&P Capital IQ (Apr 16, 2026) · comp-trading-multiples-summary.md · "
            "precedent-transactions.md · secondary-market-pricing.md",
            size=8, color=GRAY)

# =================================================================
# SLIDE 2: STRIP ONLY (drop-in element)
# =================================================================
s2 = p.slides.add_slide(LAYOUT)
add_textbox(s2, 0.45, 0.30, 19, 0.32,
            "TRIANGULATION STRIP — DROP-IN ELEMENT", size=11, bold=True, color=ORANGE)
add_textbox(s2, 0.45, 0.60, 19, 0.30,
            "Group the shapes below and paste over the bottom of the live comp slide.",
            size=10, italic=True, color=GRAY)
build_strip(s2, x0=0.45, y0=1.20, w=19.10, h=3.30)

# Save
OUT.parent.mkdir(parents=True, exist_ok=True)
p.save(str(OUT))
print(f"Wrote {OUT}")
