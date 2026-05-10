#!/usr/bin/env python3
"""Extract every slide's content from the WHOOP pptx for audit purposes."""
from pptx import Presentation
from pptx.util import Emu
import json
import sys

PPTX = "/home/user/whoop_val/Whoop Valuation Master.pptx"

prs = Presentation(PPTX)
print(f"TOTAL SLIDES: {len(prs.slides)}")
print("=" * 80)

for idx, slide in enumerate(prs.slides, start=1):
    print(f"\n{'#' * 80}")
    print(f"SLIDE {idx}")
    print(f"{'#' * 80}")
    # title
    title = None
    if slide.shapes.title is not None:
        title = slide.shapes.title.text
    print(f"TITLE: {title!r}")

    for s_idx, shape in enumerate(slide.shapes):
        try:
            left = Emu(shape.left).inches if shape.left is not None else None
            top = Emu(shape.top).inches if shape.top is not None else None
            width = Emu(shape.width).inches if shape.width is not None else None
            height = Emu(shape.height).inches if shape.height is not None else None
        except Exception:
            left = top = width = height = None

        kind = shape.shape_type
        name = shape.name
        print(f"\n  [Shape {s_idx}] type={kind} name={name!r} pos=({left},{top}) size=({width},{height})")

        # text frame
        if shape.has_text_frame:
            tf = shape.text_frame
            for p_idx, para in enumerate(tf.paragraphs):
                txt = para.text
                if txt.strip():
                    print(f"    P{p_idx}: {txt}")

        # table
        if shape.has_table:
            print(f"    TABLE rows={len(shape.table.rows)} cols={len(shape.table.columns)}")
            for r_idx, row in enumerate(shape.table.rows):
                for c_idx, cell in enumerate(row.cells):
                    ctxt = cell.text.strip().replace("\n", " | ")
                    if ctxt:
                        print(f"      [{r_idx},{c_idx}] {ctxt}")

        # chart
        if shape.has_chart:
            chart = shape.chart
            try:
                ctype = chart.chart_type
            except Exception:
                ctype = "?"
            print(f"    CHART type={ctype}")
            try:
                for plot_idx, plot in enumerate(chart.plots):
                    cats = list(plot.categories)
                    print(f"      plot {plot_idx} categories={cats}")
                    for s_i, series in enumerate(plot.series):
                        try:
                            vals = list(series.values)
                        except Exception:
                            vals = "<unreadable>"
                        try:
                            sname = series.name
                        except Exception:
                            sname = "?"
                        print(f"      plot {plot_idx} series {s_i} name={sname!r} vals={vals}")
            except Exception as e:
                print(f"      CHART data unreadable: {e}")

        # picture
        if kind == 13:  # PICTURE
            try:
                img = shape.image
                fname = getattr(img, "filename", None)
                ext = getattr(img, "ext", None)
                print(f"    PICTURE filename={fname} ext={ext} size={len(img.blob)} bytes")
            except Exception as e:
                print(f"    PICTURE error: {e}")

        # group? recurse
        if kind == 6:  # GROUP
            print("    [GROUP - members below]")
            for g_idx, g_shape in enumerate(shape.shapes):
                if g_shape.has_text_frame:
                    for p in g_shape.text_frame.paragraphs:
                        if p.text.strip():
                            print(f"      G{g_idx}: {p.text}")
                if g_shape.has_table:
                    print(f"      G{g_idx} TABLE rows={len(g_shape.table.rows)} cols={len(g_shape.table.columns)}")
                    for r_idx, row in enumerate(g_shape.table.rows):
                        for c_idx, cell in enumerate(row.cells):
                            ctxt = cell.text.strip().replace("\n", " | ")
                            if ctxt:
                                print(f"        [{r_idx},{c_idx}] {ctxt}")
                if g_shape.has_chart:
                    try:
                        gctype = g_shape.chart.chart_type
                    except Exception:
                        gctype = "?"
                    print(f"      G{g_idx} CHART type={gctype}")
